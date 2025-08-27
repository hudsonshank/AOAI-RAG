import json
import logging
import requests
import azure.functions as func
from azure.storage.blob import BlobServiceClient
from azure.keyvault.secrets import SecretClient
from azure.identity import DefaultAzureCredential
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
import os
import datetime
from typing import Dict, Any, List
import hashlib
import random
from pptx import Presentation
from io import BytesIO
import openpyxl
from urllib.parse import quote
import re

# Create the function app instance
app = func.FunctionApp()

class ClientMetadataExtractor:
    """Extract client metadata from Autobahn Consultants SharePoint paths"""
    
    def __init__(self):
        # Pattern for client folders: /ClientName (PM-X)/
        self.client_pattern = re.compile(r'/([^/]+) \(PM-([A-Z])\)/')
        
        # PM name mappings
        self.pm_names = {
            'C': 'Caleb',
            'S': 'Sam', 
            'K': 'Katherine'
        }
        
        # Internal/tool folders (non-client specific)
        self.internal_folders = {
            'autobahn tools',
            'autobook and articles', 
            'pm & apm training materials',
            'archive [do not use]',
            'rock the recession',
            'to sort',
            'parking lot',
            'jennifur',
            'general resources',
            'the glovebox',
            'operations'
        }
        
        # Document category mappings based on common folder structures
        self.category_keywords = {
            'financials': ['financial', 'finance', 'budget', 'accounting', 'statements'],
            'meetings': ['meeting', 'notes', 'agenda', 'minutes', 'quarterly'],
            'projects': ['project', 'initiative', 'implementation'],
            'documents': ['handout', 'template', 'form', 'policy'],
            'behavioral': ['behavioral', 'profile', 'assessment', 'culture'],
            'training': ['training', 'onboarding', 'handbook'],
            'archived': ['archive', 'archived', 'old'],
            'airplane': ['airplane', 'travel', 'dashboard']
        }
    
    def extract_client_info(self, document_path: str) -> Dict[str, Any]:
        """Extract client information from document path"""
        if not document_path:
            return self._create_unknown_client_info(document_path)
            
        # Check for client pattern
        client_match = self.client_pattern.search(document_path)
        
        if client_match:
            client_name = client_match.group(1).strip()
            pm_initial = client_match.group(2)
            pm_name = self.pm_names.get(pm_initial, "Unknown")
            
            # Extract folder path (everything after client folder)
            client_folder_end = client_match.end()
            remaining_path = document_path[client_folder_end:] if client_folder_end < len(document_path) else ""
            
            # Determine document category
            category = self._determine_document_category(document_path)
            
            return {
                "client_name": client_name,
                "pm_initial": pm_initial,
                "pm_name": pm_name,
                "folder_path": remaining_path,
                "is_client_specific": True,
                "document_category": category
            }
        else:
            # Check if it's an internal/tool folder
            path_lower = document_path.lower()
            for internal_folder in self.internal_folders:
                if f'/{internal_folder}/' in path_lower or path_lower.startswith(f'/{internal_folder}/'):
                    category = self._determine_document_category(document_path)
                    return {
                        "client_name": "Autobahn Internal",
                        "pm_initial": "N/A",
                        "pm_name": "N/A",
                        "folder_path": document_path,
                        "is_client_specific": False,
                        "document_category": category
                    }
            
            # Unknown/uncategorized
            return self._create_unknown_client_info(document_path)
    
    def _create_unknown_client_info(self, document_path: str) -> Dict[str, Any]:
        """Create client info for unknown/uncategorized documents"""
        category = self._determine_document_category(document_path)
        return {
            "client_name": "Uncategorized",
            "pm_initial": "N/A",
            "pm_name": "N/A", 
            "folder_path": document_path,
            "is_client_specific": False,
            "document_category": category
        }
    
    def _determine_document_category(self, document_path: str) -> str:
        """Determine document category based on path keywords"""
        path_lower = document_path.lower()
        
        # Check each category
        for category, keywords in self.category_keywords.items():
            if any(keyword in path_lower for keyword in keywords):
                return category
                
        # Default category
        return "general"

class DocumentChunker:
    """Chunk documents optimally for RAG"""
    
    def __init__(self):
        self.chunk_size = 1000
        self.chunk_overlap = 200
        self.min_chunk_size = 100
    
    def chunk_document(self, content: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Split document into RAG-optimized chunks"""
        if not content or len(content.strip()) < self.min_chunk_size:
            return []
        
        chunks = []
        content = content.strip()
        
        # Simple sentence-aware chunking
        sentences = re.split(r'(?<=[.!?])\s+', content)
        
        current_chunk = ""
        chunk_index = 0
        
        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk) + len(sentence) > self.chunk_size and current_chunk:
                # Save current chunk
                chunks.append({
                    "chunk": current_chunk.strip(),
                    "chunk_id": f"{metadata['document_id']}_{chunk_index}",
                    "parent_id": metadata['document_id'],
                    "chunk_index": chunk_index,
                    **metadata
                })
                
                # Start new chunk with overlap
                if len(current_chunk) > self.chunk_overlap:
                    overlap_start = len(current_chunk) - self.chunk_overlap
                    current_chunk = current_chunk[overlap_start:] + " " + sentence
                else:
                    current_chunk = sentence
                
                chunk_index += 1
            else:
                current_chunk += " " + sentence if current_chunk else sentence
        
        # Add final chunk if it exists
        if current_chunk.strip():
            chunks.append({
                "chunk": current_chunk.strip(),
                "chunk_id": f"{metadata['document_id']}_{chunk_index}",
                "parent_id": metadata['document_id'],
                "chunk_index": chunk_index,
                **metadata
            })
        
        return chunks

    def chunk_excel_document(self, excel_data: Dict[str, Any], metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Create separate chunks for each Excel sheet"""
        if excel_data.get("type") != "excel_sheets":
            logging.warning(f"Invalid Excel data type: {excel_data.get('type')}")
            return []
        
        chunks = []
        sheet_index = 0
        
        for sheet_name, sheet_content in excel_data["sheets"].items():
            if not sheet_content or len(sheet_content.strip()) < self.min_chunk_size:
                continue
            
            # Create metadata for this sheet
            sheet_metadata = {
                **metadata,
                "sheet_name": sheet_name,
                "sheet_index": sheet_index,
                "is_excel_sheet": True,
                "excel_sheet_count": excel_data["total_sheets"]
            }
            
            # Create chunk for this sheet
            chunk = {
                "chunk": sheet_content.strip(),
                "chunk_id": f"{metadata['document_id']}_sheet_{sheet_index}",
                "parent_id": metadata['document_id'],
                "chunk_index": sheet_index,
                "chunk_type": "excel_sheet",
                **sheet_metadata
            }
            
            chunks.append(chunk)
            sheet_index += 1
        
        logging.info(f"Created {len(chunks)} sheet-based chunks for Excel file")
        return chunks

@app.route(route="ProcessSingleDocument", methods=["POST"])
def process_single_document(req: func.HttpRequest) -> func.HttpResponse:
    """
    Enhanced HTTP endpoint for processing documents with client metadata and RAG optimization
    
    Expected JSON body:
    For folder processing:
    {
        "site_name": "Clients",
        "folder_paths": [
            "Shared Documents/ClientName (PM-C)/Financials",
            "Shared Documents/AnotherClient (PM-S)/Meetings"
        ]
    }
    
    For single file processing:
    {
        "site_name": "Clients",
        "folder_path": "Shared Documents/ClientName (PM-C)/Financials",
        "file_name": "document.pdf"
    }
    """
    logging.info('Enhanced document processing with client metadata triggered')
    try:
        req_body = req.get_json()
        if not req_body or not req_body.get('site_name'):
            return func.HttpResponse(
                json.dumps({"error": "Provide site_name in request body"}),
                status_code=400,
                mimetype="application/json"
            )
        
        site_name = req_body['site_name']
        file_name = req_body.get('file_name')
        
        # Check if processing a single file
        if file_name:
            folder_path = req_body.get('folder_path')
            processor = EnhancedDocumentProcessor()
            result = processor.process_single_file(site_name, folder_path, file_name)
            return func.HttpResponse(
                json.dumps(result, indent=2),
                status_code=200,
                mimetype="application/json"
            )
        
        # Support both single and multiple folder paths
        folder_paths = req_body.get('folder_paths')
        if not folder_paths:
            folder_path = req_body.get('folder_path')
            if not folder_path:
                return func.HttpResponse(
                    json.dumps({"error": "Provide either folder_paths or folder_path"}),
                    status_code=400,
                    mimetype="application/json"
                )
            folder_paths = [folder_path]
        
        processor = EnhancedDocumentProcessor()
        all_results = []
        
        for folder_path in folder_paths:
            logging.info(f'Processing folder: {folder_path}')
            result = processor.process_site_documents(site_name, folder_path)
            all_results.append({
                "folder_path": folder_path,
                "result": result
            })
        
        return func.HttpResponse(
            json.dumps({
                "processed_folders": len(folder_paths),
                "results": all_results
            }, indent=2),
            status_code=200,
            mimetype="application/json"
        )
        
    except Exception as e:
        logging.error(f'Error in process_single_document: {str(e)}')
        return func.HttpResponse(
            json.dumps({"error": str(e)}),
            status_code=500,
            mimetype="application/json"
        )

class EnhancedDocumentProcessor:
    def __init__(self):
        """Initialize Azure services and enhanced processing components"""
        try:
            # Load environment variables
            self.storage_connection = os.environ.get('AZURE_STORAGE_CONNECTION_STRING')
            self.key_vault_url = os.environ.get('AZURE_KEY_VAULT_URL')
            self.doc_intelligence_endpoint = os.environ.get('DOCUMENT_INTELLIGENCE_ENDPOINT')
            self.doc_intelligence_key = os.environ.get('DOCUMENT_INTELLIGENCE_KEY')
            self.pii_function_url = os.environ.get('PII_FUNCTION_URL')
            self.pii_function_key = os.environ.get('PII_FUNCTION_KEY')
            
            # Cost Control Settings
            self.test_mode = os.environ.get('TEST_MODE', 'false').lower() == 'true'
            self.max_documents_per_run = int(os.environ.get('MAX_DOCUMENTS_PER_RUN', '50'))
            self.max_file_size_mb = int(os.environ.get('MAX_FILE_SIZE_MB', '100'))
            self.skip_doc_intelligence_for_large_files = True
            
            # Initialize Azure clients
            self.storage_client = BlobServiceClient.from_connection_string(self.storage_connection)
            self.key_vault_client = SecretClient(vault_url=self.key_vault_url, credential=DefaultAzureCredential())
            self.doc_intelligence_client = DocumentAnalysisClient(
                endpoint=self.doc_intelligence_endpoint,
                credential=AzureKeyCredential(self.doc_intelligence_key)
            )
            
            # Initialize enhanced components
            self.client_extractor = ClientMetadataExtractor()
            self.chunker = DocumentChunker()
            
            # Get Graph API token
            self.graph_token = self._get_graph_token()
            
            # Define supported file types
            self.supported_extensions = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt'}
            self.lightweight_extensions = {'.txt', '.docx', '.doc'}
            
            # SharePoint sites to scan
            self.sharepoint_sites = ["APM Board", "Clients", "Autobahn Overdrive Artificial Intelligence", "PM", "Operations", "The Glovebox", "Communication site"]
            
            # Enhanced folder filtering
            self.skip_folders = [
                "template emails", "archive [do not use]", "email templates", 
                "6. templates", "/templates/", "ser00", "health nucleus",
                "recycle bin", "forms/templates", "_private", "workflow history"
            ]
            
            # Processing statistics
            self.processing_stats = {
                "documents_processed_this_run": 0,
                "total_cost_estimate": 0.0,
                "processing_start_time": datetime.datetime.utcnow()
            }
            
            logging.info('Enhanced DocumentProcessor initialized with client metadata support')
            
        except Exception as e:
            logging.error(f'Error initializing EnhancedDocumentProcessor: {str(e)}')
            raise

    def _get_graph_token(self):
        """Get Microsoft Graph API token from Key Vault"""
        try:
            token_secret = self.key_vault_client.get_secret("graph-api-token")
            return token_secret.value
        except Exception as e:
            logging.error(f'Error getting Graph token: {str(e)}')
            raise

    def process_single_file(self, site_name: str, folder_path: str, file_name: str) -> Dict[str, Any]:
        """Process a single file with enhanced metadata"""
        logging.info(f'Processing single file: {file_name} in {folder_path}')
        
        try:
            # Get file details from SharePoint
            file_info = self._get_file_info(site_name, folder_path, file_name)
            if not file_info:
                return {"error": f"File not found: {file_name}"}
            
            # Process the document
            result = self._process_single_document_enhanced(file_info, folder_path)
            
            return {
                "file_name": file_name,
                "folder_path": folder_path,
                "processing_result": result
            }
            
        except Exception as e:
            logging.error(f'Error processing single file {file_name}: {str(e)}')
            return {"error": str(e)}

    def process_site_documents(self, site_name: str, folder_path: str = None) -> Dict[str, Any]:
        """Process documents from a specific SharePoint site and folder"""
        logging.info(f'Processing documents from site: {site_name}, folder: {folder_path}')
        
        try:
            site_id = self._get_site_id(site_name)
            if not site_id:
                return {"error": f"Site not found: {site_name}"}
            
            documents = self._get_documents_from_folder(site_id, folder_path)
            
            results = {
                "site_name": site_name,
                "folder_path": folder_path,
                "total_documents_found": len(documents),
                "processed_documents": [],
                "skipped_documents": [],
                "errors": []
            }
            
            for doc in documents:
                if self.processing_stats["documents_processed_this_run"] >= self.max_documents_per_run:
                    logging.info(f'Reached max documents limit ({self.max_documents_per_run}), stopping')
                    break
                
                try:
                    result = self._process_single_document_enhanced(doc, folder_path)
                    
                    if result.get("action") == "processed":
                        results["processed_documents"].append(result)
                        self.processing_stats["documents_processed_this_run"] += 1
                    else:
                        results["skipped_documents"].append(result)
                        
                except Exception as e:
                    error_result = {
                        "document": doc.get('name', 'unknown'),
                        "error": str(e)
                    }
                    results["errors"].append(error_result)
                    logging.error(f'Error processing document {doc.get("name")}: {str(e)}')
            
            return results
            
        except Exception as e:
            logging.error(f'Error processing site documents: {str(e)}')
            return {"error": str(e)}

    def _process_single_document_enhanced(self, doc: Dict[str, Any], base_folder_path: str = None) -> Dict[str, Any]:
        """Process a single document with enhanced client metadata and RAG optimization"""
        
        doc_id = doc['id']
        doc_name = doc['name']
        doc_path = doc.get('path', doc_name)
        doc_extension = doc.get('extension', 'unknown')
        doc_size = doc.get('size', 0)
        
        # Build full document path for client metadata extraction
        full_doc_path = f"/{base_folder_path}/{doc_name}" if base_folder_path else f"/{doc_name}"
        
        processing_start = datetime.datetime.utcnow()
        
        try:
            # Check if already processed
            if self._is_document_processed(doc_id):
                logging.info(f'Document {doc_path} already processed, skipping')
                return {
                    "action": "skipped", 
                    "reason": "already_processed", 
                    "path": doc_path,
                    "extension": doc_extension
                }
            
            # Size-based processing limits
            max_size_bytes = self.max_file_size_mb * 1024 * 1024
            if doc_size > max_size_bytes:
                if self.skip_doc_intelligence_for_large_files:
                    logging.info(f'Large file {doc_path} ({doc_size:,} bytes), using fallback processing')
                    return self._process_large_file_with_fallback(doc, doc_size, full_doc_path)
                else:
                    logging.info(f'File {doc_path} exceeds size limit ({doc_size:,} bytes), skipping')
                    return {
                        "action": "skipped_size_limit",
                        "reason": f"file_too_large_{doc_size}_bytes",
                        "path": doc_path,
                        "extension": doc_extension
                    }
            
            logging.info(f'Processing document: {doc_path} ({doc_extension}, {doc_size:,} bytes)')
            
            # Download document
            doc_content = self._download_document(doc['download_url'])
            
            # Extract text content
            extracted_content = self._extract_text_content(doc_content, doc_extension, doc_name)
            
            # Extract client metadata
            client_info = self.client_extractor.extract_client_info(full_doc_path)
            
            # Handle Excel files differently - they return structured data
            if doc_extension.lower() in ['.xlsx', '.xls'] and isinstance(extracted_content, dict):
                if extracted_content.get("type") == "excel_sheets":
                    # Use Excel-specific chunking for sheet-based processing
                    chunks = self.chunker.chunk_excel_document(extracted_content, {
                        "id": doc_id,
                        "name": doc_name,
                        "document_path": full_doc_path,
                        "filename": doc_name,
                        "size": doc_size,
                        "extension": doc_extension,
                        "processed_timestamp": datetime.datetime.utcnow().isoformat(),
                        "status": "ready_for_rag",
                        "client_name": client_info["client_name"],
                        "pm_initial": client_info["pm_initial"],
                        "pm_name": client_info["pm_name"],
                        "document_category": client_info["document_category"],
                        "is_client_specific": client_info["is_client_specific"],
                        "processing_method": "excel_sheet_extraction",
                        "folder_depth": len([p for p in full_doc_path.split('/') if p]),
                        "file_extension": doc_extension.lower()
                    })
                    
                    # Calculate content length from all sheets
                    total_content_length = sum(len(content) for content in extracted_content["sheets"].values())
                    enhanced_metadata = {
                        "id": doc_id,
                        "name": doc_name,
                        "document_path": full_doc_path,
                        "filename": doc_name,
                        "size": doc_size,
                        "extension": doc_extension,
                        "processed_timestamp": datetime.datetime.utcnow().isoformat(),
                        "status": "ready_for_rag",
                        "client_name": client_info["client_name"],
                        "pm_initial": client_info["pm_initial"],
                        "pm_name": client_info["pm_name"],
                        "document_category": client_info["document_category"],
                        "is_client_specific": client_info["is_client_specific"],
                        "content_length": total_content_length,
                        "word_count": sum(len(content.split()) for content in extracted_content["sheets"].values()),
                        "character_count": total_content_length,
                        "processing_method": "excel_sheet_extraction",
                        "folder_depth": len([p for p in full_doc_path.split('/') if p]),
                        "file_extension": doc_extension.lower(),
                        "excel_sheet_count": extracted_content["total_sheets"]
                    }
                else:
                    # Excel extraction failed, skip
                    return {
                        "action": "skipped",
                        "reason": "excel_extraction_failed",
                        "path": doc_path,
                        "extension": doc_extension,
                        "error": extracted_content.get("error", "Unknown Excel error")
                    }
            else:
                # Handle regular text-based documents
                if not extracted_content or len(extracted_content.strip()) < 50:
                    return {
                        "action": "skipped",
                        "reason": "no_extractable_content",
                        "path": doc_path,
                        "extension": doc_extension
                    }
                
                # Create enhanced metadata for regular documents
                enhanced_metadata = {
                    "id": doc_id,
                    "name": doc_name,
                    "document_path": full_doc_path,
                    "filename": doc_name,
                    "size": doc_size,
                    "extension": doc_extension,
                    "processed_timestamp": datetime.datetime.utcnow().isoformat(),
                    "status": "ready_for_rag",
                    "client_name": client_info["client_name"],
                    "pm_initial": client_info["pm_initial"],
                    "pm_name": client_info["pm_name"],
                    "document_category": client_info["document_category"],
                    "is_client_specific": client_info["is_client_specific"],
                    "content_length": len(extracted_content),
                    "word_count": len(extracted_content.split()),
                    "character_count": len(extracted_content),
                    "processing_method": "azure_document_intelligence",
                    "folder_depth": len([p for p in full_doc_path.split('/') if p]),
                    "file_extension": doc_extension.lower()
                }
                
                # Chunk document for RAG using regular chunking
                chunks = self.chunker.chunk_document(extracted_content, enhanced_metadata)
            
            if not chunks:
                return {
                    "action": "skipped",
                    "reason": "no_valid_chunks",
                    "path": doc_path,
                    "extension": doc_extension
                }
            
            # Store each chunk as a separate document for RAG
            stored_chunks = 0
            for chunk in chunks:
                try:
                    self._store_rag_optimized_document(chunk)
                    stored_chunks += 1
                except Exception as e:
                    logging.error(f'Error storing chunk {chunk["chunk_id"]}: {str(e)}')
            
            processing_time = (datetime.datetime.utcnow() - processing_start).total_seconds()
            
            # Calculate content length based on document type
            if doc_extension.lower() in ['.xlsx', '.xls'] and isinstance(extracted_content, dict):
                if extracted_content.get("type") == "excel_sheets":
                    content_length = sum(len(content) for content in extracted_content["sheets"].values())
                else:
                    content_length = 0
            else:
                content_length = len(extracted_content) if extracted_content else 0
            
            result = {
                "action": "processed",
                "path": doc_path,
                "extension": doc_extension,
                "client_name": client_info["client_name"],
                "pm_name": client_info["pm_name"],
                "document_category": client_info["document_category"],
                "chunks_created": len(chunks),
                "chunks_stored": stored_chunks,
                "content_length": content_length,
                "processing_time_seconds": processing_time,
                "processing_method": enhanced_metadata.get("processing_method", "standard")
            }
            
            # Add Excel-specific information if applicable
            if doc_extension.lower() in ['.xlsx', '.xls'] and isinstance(extracted_content, dict):
                if extracted_content.get("type") == "excel_sheets":
                    result.update({
                        "excel_sheets_processed": extracted_content["total_sheets"],
                        "sheet_names": extracted_content["sheet_names"]
                    })
            
            logging.info(f'Successfully processed {doc_name}: {stored_chunks} chunks stored')
            return result
            
        except Exception as e:
            processing_time = (datetime.datetime.utcnow() - processing_start).total_seconds()
            error_result = {
                "action": "error",
                "path": doc_path,
                "extension": doc_extension,
                "error": str(e),
                "processing_time_seconds": processing_time
            }
            logging.error(f'Error processing document {doc_name}: {str(e)}')
            return error_result

    def _store_rag_optimized_document(self, chunk_data: Dict[str, Any]):
        """Store RAG-optimized document chunk in jennifur-processed container"""
        try:
            # Create document structure for RAG indexing
            rag_document = {
                "chunk_id": chunk_data["chunk_id"],
                "parent_id": chunk_data["parent_id"],
                "document_id": chunk_data["id"],
                "document_path": chunk_data["document_path"],
                "filename": chunk_data["filename"],
                "title": chunk_data["filename"],
                "chunk": chunk_data["chunk"],
                "chunk_index": chunk_data["chunk_index"],
                
                # Client metadata for filtering
                "client_name": chunk_data["client_name"],
                "pm_initial": chunk_data["pm_initial"],
                "pm_name": chunk_data["pm_name"],
                "document_category": chunk_data["document_category"],
                "is_client_specific": chunk_data["is_client_specific"],
                
                # Content metadata
                "content_length": len(chunk_data["chunk"]),
                "word_count": len(chunk_data["chunk"].split()),
                "processed_timestamp": chunk_data["processed_timestamp"],
                "status": "ready_for_rag",
                
                # Additional metadata for search
                "file_extension": chunk_data["file_extension"],
                "folder_depth": chunk_data["folder_depth"],
                "processing_method": chunk_data.get("processing_method", "enhanced"),
                
                # Excel-specific metadata (if applicable)
                "chunk_type": chunk_data.get("chunk_type", "standard"),
                "sheet_name": chunk_data.get("sheet_name"),
                "sheet_index": chunk_data.get("sheet_index"),
                "is_excel_sheet": chunk_data.get("is_excel_sheet", False),
                "excel_sheet_count": chunk_data.get("excel_sheet_count"),
                
                # Full metadata for reference
                "metadata": {
                    "size": chunk_data["size"],
                    "extension": chunk_data["extension"],
                    "processing_timestamp": chunk_data["processed_timestamp"]
                }
            }
            
            # Store in jennifur-processed container
            blob_client = self.storage_client.get_blob_client(
                container="jennifur-processed",
                blob=f"{chunk_data['chunk_id']}.json"
            )
            
            blob_client.upload_blob(
                json.dumps(rag_document, indent=2),
                overwrite=True
            )
            
            logging.info(f'RAG-optimized chunk {chunk_data["chunk_id"]} stored in jennifur-processed')
            
        except Exception as e:
            logging.error(f'Error storing RAG document: {str(e)}')
            raise

    def _is_document_processed(self, doc_id: str) -> bool:
        """Check if document has already been processed in jennifur-processed"""
        try:
            # Check for any chunks from this parent document
            container_client = self.storage_client.get_container_client("jennifur-processed")
            blobs = container_client.list_blobs(name_starts_with=doc_id)
            return any(True for _ in blobs)
        except Exception:
            return False

    # Include other necessary methods from the original file
    def _get_graph_token(self):
        """Get Microsoft Graph API token from Key Vault"""
        try:
            token_secret = self.key_vault_client.get_secret("graph-api-token")
            return token_secret.value
        except Exception as e:
            logging.error(f'Error getting Graph token: {str(e)}')
            raise

    def _get_site_id(self, site_name: str) -> str:
        """Get SharePoint site ID"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            url = f"https://graph.microsoft.com/v1.0/sites?search={quote(site_name)}"
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            
            sites = response.json().get('value', [])
            for site in sites:
                if site['displayName'] == site_name:
                    return site['id']
            return None
        except Exception as e:
            logging.error(f'Error getting site ID for {site_name}: {str(e)}')
            return None

    def _get_documents_from_folder(self, site_id: str, folder_path: str = None) -> List[Dict[str, Any]]:
        """Get documents from SharePoint folder"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            
            if folder_path:
                url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/root:/{folder_path}:/children"
            else:
                url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/root/children"
            
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            
            items = response.json().get('value', [])
            documents = []
            
            for item in items:
                if 'file' in item:  # It's a file, not a folder
                    file_extension = os.path.splitext(item['name'])[1].lower()
                    if file_extension in self.supported_extensions:
                        documents.append({
                            'id': item['id'],
                            'name': item['name'],
                            'path': item.get('parentReference', {}).get('path', '') + '/' + item['name'],
                            'extension': file_extension,
                            'size': item.get('size', 0),
                            'download_url': item.get('@microsoft.graph.downloadUrl')
                        })
            
            return documents
            
        except Exception as e:
            logging.error(f'Error getting documents from folder: {str(e)}')
            return []

    def _download_document(self, download_url: str) -> bytes:
        """Download document content from OneDrive/SharePoint"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            response = requests.get(download_url, headers=headers)
            response.raise_for_status()
            return response.content
            
        except Exception as e:
            logging.error(f'Error downloading document: {str(e)}')
            raise

    def _extract_text_content(self, doc_content: bytes, extension: str, filename: str) -> str:
        """Extract text content from document using appropriate method"""
        try:
            if extension.lower() in ['.pdf', '.docx', '.doc']:
                # Use Azure Document Intelligence for optimal extraction
                return self._extract_with_document_intelligence(doc_content, extension)
            elif extension.lower() in ['.txt']:
                return doc_content.decode('utf-8', errors='ignore')
            elif extension.lower() in ['.xlsx', '.xls']:
                return self._extract_from_excel(doc_content)
            elif extension.lower() in ['.pptx', '.ppt']:
                return self._extract_from_powerpoint(doc_content)
            else:
                return self._extract_with_document_intelligence(doc_content, extension)
                
        except Exception as e:
            logging.error(f'Error extracting text from {filename}: {str(e)}')
            return ""

    def _extract_with_document_intelligence(self, doc_content: bytes, extension: str) -> str:
        """Extract text using Azure Document Intelligence"""
        try:
            poller = self.doc_intelligence_client.begin_analyze_document(
                "prebuilt-read",
                document=BytesIO(doc_content)
            )
            result = poller.result()
            
            return result.content
            
        except Exception as e:
            logging.error(f'Document Intelligence extraction failed: {str(e)}')
            return ""

    def _extract_from_excel(self, doc_content: bytes) -> Dict[str, Any]:
        """Extract text from Excel files with visible text and sheet-based structure"""
        try:
            workbook = openpyxl.load_workbook(BytesIO(doc_content), read_only=True, data_only=True)
            sheets_data = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = []
                
                # Add sheet header
                sheet_content.append(f"Sheet: {sheet_name}")
                
                # Process each row
                for row_num, row in enumerate(sheet.iter_rows(), 1):
                    row_cells = []
                    
                    for col_num, cell in enumerate(row, 1):
                        if cell is not None:
                            # Get the displayed value (not formula)
                            cell_value = cell.value
                            
                            # Handle different data types
                            if cell_value is not None:
                                if isinstance(cell_value, (int, float)):
                                    # Format numbers appropriately
                                    if isinstance(cell_value, int):
                                        cell_text = str(cell_value)
                                    else:
                                        # Format floats to avoid excessive decimals
                                        cell_text = f"{cell_value:.2f}".rstrip('0').rstrip('.')
                                elif isinstance(cell_value, datetime.datetime):
                                    # Format dates consistently
                                    cell_text = cell_value.strftime("%Y-%m-%d %H:%M:%S")
                                elif isinstance(cell_value, datetime.date):
                                    cell_text = cell_value.strftime("%Y-%m-%d")
                                elif isinstance(cell_value, datetime.time):
                                    cell_text = cell_value.strftime("%H:%M:%S")
                                else:
                                    cell_text = str(cell_value).strip()
                                
                                # Only add non-empty cells
                                if cell_text and cell_text != "None":
                                    row_cells.append(cell_text)
                    
                    # Only add rows with content
                    if row_cells:
                        row_text = " | ".join(row_cells)
                        sheet_content.append(f"Row {row_num}: {row_text}")
                
                # Store sheet content
                if len(sheet_content) > 1:  # More than just the sheet header
                    sheets_data[sheet_name] = "\n".join(sheet_content)
            
            workbook.close()
            
            # Return structured data for sheet-based chunking
            return {
                "type": "excel_sheets",
                "sheets": sheets_data,
                "total_sheets": len(sheets_data),
                "sheet_names": list(sheets_data.keys())
            }
            
        except Exception as e:
            logging.error(f'Excel extraction failed: {str(e)}')
            return {
                "type": "excel_error",
                "error": str(e),
                "sheets": {},
                "total_sheets": 0,
                "sheet_names": []
            }

    def _extract_from_powerpoint(self, doc_content: bytes) -> str:
        """Extract text from PowerPoint files"""
        try:
            presentation = Presentation(BytesIO(doc_content))
            content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                content.append(f"Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        content.append(shape.text)
            
            return "\n".join(content)
            
        except Exception as e:
            logging.error(f'PowerPoint extraction failed: {str(e)}')
            return ""

    def _get_file_info(self, site_name: str, folder_path: str, file_name: str) -> Dict[str, Any]:
        """Get specific file information from SharePoint"""
        try:
            site_id = self._get_site_id(site_name)
            if not site_id:
                return None
            
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/root:/{folder_path}/{file_name}"
            
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            
            item = response.json()
            
            return {
                'id': item['id'],
                'name': item['name'],
                'path': item.get('parentReference', {}).get('path', '') + '/' + item['name'],
                'extension': os.path.splitext(item['name'])[1].lower(),
                'size': item.get('size', 0),
                'download_url': item.get('@microsoft.graph.downloadUrl')
            }
            
        except Exception as e:
            logging.error(f'Error getting file info for {file_name}: {str(e)}')
            return None

    def _process_large_file_with_fallback(self, doc: Dict[str, Any], doc_size: int, full_doc_path: str) -> Dict[str, Any]:
        """Process large files with fallback method"""
        try:
            doc_content = self._download_document(doc['download_url'])
            
            # Simple text extraction for large files
            if doc['extension'].lower() == '.txt':
                content = doc_content.decode('utf-8', errors='ignore')
            else:
                content = f"Large file detected: {doc['name']} ({doc_size:,} bytes). Content extraction limited for cost control."
            
            # Extract client metadata
            client_info = self.client_extractor.extract_client_info(full_doc_path)
            
            # Create single chunk for large files
            chunk_data = {
                "chunk_id": f"{doc['id']}_0",
                "parent_id": doc['id'],
                "chunk": content[:2000],  # Limit content for large files
                "chunk_index": 0,
                "id": doc['id'],
                "name": doc['name'],
                "document_path": full_doc_path,
                "filename": doc['name'],
                "size": doc_size,
                "extension": doc['extension'],
                "processed_timestamp": datetime.datetime.utcnow().isoformat(),
                **client_info,
                "content_length": len(content),
                "word_count": len(content.split()),
                "character_count": len(content),
                "processing_method": "fallback_large_file",
                "folder_depth": len([p for p in full_doc_path.split('/') if p]),
                "file_extension": doc['extension'].lower()
            }
            
            self._store_rag_optimized_document(chunk_data)
            
            return {
                "action": "processed_fallback",
                "path": doc.get('path', doc['name']),
                "extension": doc['extension'],
                "client_name": client_info["client_name"],
                "pm_name": client_info["pm_name"],
                "reason": "large_file_fallback",
                "size_bytes": doc_size
            }
            
        except Exception as e:
            return {
                "action": "error",
                "path": doc.get('path', doc['name']),
                "extension": doc['extension'],
                "error": str(e)
            }