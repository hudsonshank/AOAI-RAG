import json
import logging
import requests
import azure.functions as func
from azure.storage.blob import BlobServiceClient
from azure.keyvault.secrets import SecretClient
from azure.identity import DefaultAzureCredential
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
import os
import datetime
from typing import Dict, Any, List
import hashlib
import random
from pptx import Presentation  # Add this import for PowerPoint extraction
from io import BytesIO
import openpyxl
from urllib.parse import quote

def main(req: func.HttpRequest) -> func.HttpResponse:
    """
    Azure Function v2 main entry point for document processing
    """
    logging.info('Manual document processing triggered')
    
    try:
        req_body = req.get_json()
        if not req_body or not req_body.get('site_name'):
            return func.HttpResponse(
                json.dumps({"error": "Provide site_name in request body"}),
                status_code=400,
                mimetype="application/json"
            )
        
        site_name = req_body['site_name']
        file_name = req_body.get('file_name')
        
        # Check if processing a single file
        if file_name:
            folder_path = req_body.get('folder_path')
            processor = DocumentProcessor()
            result = processor.process_single_file(site_name, folder_path, file_name)
            return func.HttpResponse(
                json.dumps(result, indent=2),
                status_code=200,
                mimetype="application/json"
            )
        
        # Support both single and multiple folder paths
        folder_paths = req_body.get('folder_paths')
        if not folder_paths:
            folder_path = req_body.get('folder_path')
            folder_paths = [folder_path] if folder_path else [None]
        
        # Get recursive option (defaults to False if not specified)
        recursive = req_body.get('recursive', False)

        processor = DocumentProcessor()
        results = []
        for folder_path in folder_paths:
            if folder_path:
                result = processor.process_site_documents(site_name, folder_path, recursive)
            else:
                result = processor.process_site_documents(site_name, None, recursive)
            results.append(result)

        return func.HttpResponse(
            json.dumps(results, indent=2),
            status_code=200,
            mimetype="application/json"
        )
    except Exception as e:
        logging.error(f'Error processing document: {str(e)}')
        return func.HttpResponse(
            json.dumps({"error": str(e)}),
            status_code=500,
            mimetype="application/json"
        )

class DocumentProcessor:
    def __init__(self):
        """Initialize Azure services and Graph API connection"""
        try:
            # Load environment variables
            self.storage_connection = os.environ.get('AZURE_STORAGE_CONNECTION_STRING')
            self.key_vault_url = os.environ.get('AZURE_KEY_VAULT_URL')
            self.doc_intelligence_endpoint = os.environ.get('DOCUMENT_INTELLIGENCE_ENDPOINT')
            self.doc_intelligence_key = os.environ.get('DOCUMENT_INTELLIGENCE_KEY')
            self.pii_function_url = os.environ.get('PII_FUNCTION_URL')
            self.pii_function_key = os.environ.get('PII_FUNCTION_KEY')
            
            # *** Cost Control Settings ***
            self.test_mode = os.environ.get('TEST_MODE', 'false').lower() == 'true'
            self.max_documents_per_run = int(os.environ.get('MAX_DOCUMENTS_PER_RUN', '100'))  # Increased from 50
            self.max_file_size_mb = int(os.environ.get('MAX_FILE_SIZE_MB', '100'))
            self.skip_doc_intelligence_for_large_files = True
            
            # Initialize Azure clients
            self.storage_client = BlobServiceClient.from_connection_string(self.storage_connection)
            self.key_vault_client = SecretClient(vault_url=self.key_vault_url, credential=DefaultAzureCredential())
            self.doc_intelligence_client = DocumentAnalysisClient(
                endpoint=self.doc_intelligence_endpoint,
                credential=AzureKeyCredential(self.doc_intelligence_key)
            )
            
            # Get Graph API token
            self.graph_token = self._get_graph_token()
            
            # Define supported file types with size preferences
            self.supported_extensions = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt'}
            self.lightweight_extensions = {'.txt', '.docx', '.doc'}  # Process these first
            
            # SharePoint sites to scan
            self.sharepoint_sites = ["APM Board", "Clients", "Autobahn Overdrive Artificial Intelligence", "PM", "Operations", "The Glovebox", "Communication site"]
            
            # Enhanced folder filtering
            self.skip_folders = [
                "template emails", "archive [do not use]", "email templates", 
                "6. templates", "/templates/", "ser00", "health nucleus",
                "recycle bin", "forms/templates", "_private", "workflow history"
            ]
            
            # Processing statistics
            self.processing_stats = {
                "documents_processed_this_run": 0,
                "total_cost_estimate": 0.0,
                "processing_start_time": datetime.datetime.utcnow()
            }
            
            logging.info(f'DocumentProcessor initialized - Test Mode: {self.test_mode}, Max Docs: {self.max_documents_per_run}')
            
        except Exception as e:
            logging.error(f'Failed to initialize DocumentProcessor: {str(e)}')
            raise
    
    def _get_graph_token(self) -> str:
        """Get Microsoft Graph API access token"""
        try:
            client_secret = self.key_vault_client.get_secret("Jennifur-Client-Secret").value
            
            token_data = {
                'client_id': '8bb9c943-15df-4b66-b52f-4c40debdee88',  # Your Jennifur-AI-Assistant app ID
                'client_secret': client_secret,
                'scope': 'https://graph.microsoft.com/.default',
                'grant_type': 'client_credentials'
            }
            
            response = requests.post(
                'https://login.microsoftonline.com/95be67cb-cbae-432a-b9ca-d3befaae2e0e/oauth2/v2.0/token',
                data=token_data
            )
            response.raise_for_status()
            
            return response.json()['access_token']
            
        except Exception as e:
            logging.error(f'Failed to get Graph API token: {str(e)}')
            raise

    def process_site_documents(self, site_name: str, folder_path: str = None, recursive: bool = False) -> Dict[str, Any]:
        """Process documents from a specific SharePoint site with comprehensive error tracking"""
        
        site_results = {
            "site_name": site_name,
            "test_mode": self.test_mode,
            "documents_found": 0,
            "documents_processed": 0,
            "documents_quarantined": 0,
            "documents_flagged": 0,
            "documents_skipped": 0,
            "documents_failed": 0,
            "documents_skipped_size_limit": 0,
            "documents_skipped_extraction_failed": 0,
            "folders_skipped": 0,
            "processing_errors": [],
            "file_type_summary": {},
            "cost_estimate": 0.0,
            "total_processing_time_seconds": 0
        }
        
        processing_start = datetime.datetime.utcnow()
        
        try:
            # Get site ID
            site_id = self._get_site_id(site_name)
            if not site_id:
                logging.warning(f'⚠️ Site not found: {site_name}')
                return site_results
            
            # Load checkpoint for resuming from last position
            checkpoint = self._load_checkpoint(site_name, folder_path)
            
            # Get documents from site
            documents = self._get_site_documents(site_id, folder_path, site_name, checkpoint, recursive)
            site_results["documents_found"] = len(documents)
            logging.info(f"[BATCH] Found {len(documents)} documents in SharePoint for site '{site_name}' and folder '{folder_path or '/'}' (recursive: {recursive})")
            
            # Add folder path and recursive option to results for tracking
            site_results["folder_path"] = folder_path or "/"
            site_results["recursive"] = recursive

            # Prioritize documents for RAG ingestion
            documents_to_process = self._prioritize_documents_for_testing(documents)
            logging.info(f"[BATCH] Prioritized {len(documents_to_process)} documents for processing (max per run: {self.max_documents_per_run})")

            processed_count = 0
            # Process each document with cost tracking
            for doc in documents_to_process:
                # Check if we've hit our processing limit
                if self.processing_stats["documents_processed_this_run"] >= self.max_documents_per_run:
                    logging.info(f'🛑 Hit processing limit ({self.max_documents_per_run}), stopping')
                    break
                
                try:
                    result = self._process_single_document_with_cost_control(doc)
                    
                    # Update counters
                    action = result.get("action", "unknown")
                    extension = result.get("extension", "unknown")
                    cost = result.get("cost_estimate", 0.0)
                    
                    # Update processing stats
                    if action in ["processed", "quarantined", "flagged"]:
                        self.processing_stats["documents_processed_this_run"] += 1
                        self.processing_stats["total_cost_estimate"] += cost
                        site_results["cost_estimate"] += cost
                    
                    # Count by action
                    if action == "processed":
                        site_results["documents_processed"] += 1
                    elif action == "quarantined":
                        site_results["documents_quarantined"] += 1
                    elif action == "flagged":
                        site_results["documents_flagged"] += 1
                    elif action == "skipped":
                        site_results["documents_skipped"] += 1
                        if result.get("reason") == "already_processed":
                            logging.debug(f'⏭️ Document already processed: {result.get("path", "unknown")}')
                    elif action == "skipped_size_limit":
                        site_results["documents_skipped_size_limit"] += 1
                    elif action == "skipped_extraction_failed":
                        site_results["documents_skipped_extraction_failed"] += 1
                    else:
                        site_results["documents_failed"] += 1
                        site_results["processing_errors"].append({
                            "file": result.get("path", "unknown"),
                            "reason": result.get("reason", "unknown"),
                            "extension": extension
                        })
                    
                    # Count by file type
                    site_results["file_type_summary"][extension] = site_results["file_type_summary"].get(extension, 0) + 1
                    
                except Exception as e:
                    logging.error(f'❌ Critical error processing document {doc.get("name", "unknown")}: {str(e)}')
                    site_results["documents_failed"] += 1
                    site_results["processing_errors"].append({
                        "file": doc.get("path", doc.get("name", "unknown")),
                        "reason": f"critical_error: {str(e)}",
                        "extension": doc.get("extension", "unknown")
                    })
            
            # Calculate total processing time
            processing_end = datetime.datetime.utcnow()
            site_results["total_processing_time_seconds"] = round((processing_end - processing_start).total_seconds(), 2)
            
            # Enhanced logging with cost information
            logging.info(f'Site {site_name} processing complete:')
            logging.info(f'   Processed: {site_results["documents_processed"]}')
            logging.info(f'   Quarantined: {site_results["documents_quarantined"]}')
            logging.info(f'   Flagged: {site_results["documents_flagged"]}')
            logging.info(f'   Skipped (already processed): {site_results["documents_skipped"]}')
            logging.info(f'   Skipped (size limit): {site_results["documents_skipped_size_limit"]}')
            logging.info(f'   Skipped (extraction failed): {site_results["documents_skipped_extraction_failed"]}')
            logging.info(f'   Failed: {site_results["documents_failed"]}')
            logging.info(f'   Estimated cost: ${site_results["cost_estimate"]:.4f}')
            logging.info(f'   Total time: {site_results["total_processing_time_seconds"]}s')
            logging.info(f"[BATCH] Actually processed {processed_count} documents in this run.")
            
            return site_results
            
        except Exception as e:
            processing_end = datetime.datetime.utcnow()
            site_results["total_processing_time_seconds"] = round((processing_end - processing_start).total_seconds(), 2)
            error_msg = f'Error processing site {site_name}: {str(e)}'
            logging.error(error_msg)
            site_results["processing_errors"].append({
                "site": site_name,
                "reason": f"site_processing_error: {str(e)}",
                "extension": "N/A"
            })
            return site_results

    def _get_site_id(self, site_name: str) -> str:
        """Get SharePoint site ID by name"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            
            # Search for site
            response = requests.get(
                f'https://graph.microsoft.com/v1.0/sites?search={site_name}',
                headers=headers
            )
            response.raise_for_status()
            
            sites = response.json().get('value', [])
            if sites:
                return sites[0]['id']
            
            return None
            
        except Exception as e:
            logging.error(f'Error getting site ID for {site_name}: {str(e)}')
            return None

    def _get_checkpoint_blob_name(self, site_name: str, folder_path: str = None) -> str:
        """Generate checkpoint blob name for a site and folder path"""
        import re
        
        # Sanitize site name
        safe_site_name = re.sub(r'[^a-zA-Z0-9-_]', '_', site_name)
        
        # Sanitize folder path
        if folder_path:
            safe_folder_path = re.sub(r'[^a-zA-Z0-9-_/]', '_', folder_path)
            safe_folder_path = safe_folder_path.strip('/').replace('/', '_')
        else:
            safe_folder_path = 'root'
        
        return f"checkpoints/{safe_site_name}/{safe_folder_path}.json"

    def _load_checkpoint(self, site_name: str, folder_path: str = None) -> Dict[str, Any]:
        """Load checkpoint from Azure Blob Storage"""
        try:
            blob_name = self._get_checkpoint_blob_name(site_name, folder_path)
            blob_client = self.storage_client.get_blob_client(
                container="processed-documents", 
                blob=blob_name
            )
            
            if blob_client.exists():
                checkpoint_data = blob_client.download_blob().readall()
                checkpoint = json.loads(checkpoint_data.decode('utf-8'))
                logging.info(f'📍 Loaded checkpoint for {site_name}/{folder_path or "root"}')
                return checkpoint
            else:
                logging.info(f'📍 No checkpoint found for {site_name}/{folder_path or "root"}')
                return {}
                
        except Exception as e:
            logging.warning(f'Failed to load checkpoint for {site_name}/{folder_path or "root"}: {str(e)}')
            return {}
    
    def _save_checkpoint(self, site_name: str, folder_path: str = None, next_link: str = None) -> None:
        """Save checkpoint to Azure Blob Storage"""
        try:
            checkpoint = {
                "site_name": site_name,
                "folder_path": folder_path or "/",
                "last_processed_url": next_link,
                "timestamp": datetime.datetime.utcnow().isoformat()
            }
            
            blob_name = self._get_checkpoint_blob_name(site_name, folder_path)
            blob_client = self.storage_client.get_blob_client(
                container="processed-documents", 
                blob=blob_name
            )
            
            checkpoint_json = json.dumps(checkpoint, indent=2)
            blob_client.upload_blob(
                checkpoint_json, 
                overwrite=True, 
                content_type='application/json'
            )
            
            logging.info(f'💾 Saved checkpoint for {site_name}/{folder_path or "root"}')
            
        except Exception as e:
            logging.error(f'Failed to save checkpoint for {site_name}/{folder_path or "root"}: {str(e)}')

    def _get_site_documents(self, site_id: str, folder_path: str = None, site_name: str = None, checkpoint: Dict[str, Any] = None, recursive: bool = False) -> List[Dict[str, Any]]:
        """Get all documents from a SharePoint site, optionally targeting a specific folder with checkpoint support"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            documents = []
            
            # Get document libraries in the site
            drives_response = requests.get(
                f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives',
                headers=headers
            )
            drives_response.raise_for_status()
            
            drives_data = drives_response.json()
            # Log all drive names and IDs for debugging
            logging.info(f"[DEBUG] Drives for site {site_name} (site_id={site_id}):")
            for drive in drives_data['value']:
                logging.info(f"[DEBUG] Drive name: '{drive.get('name')}', ID: {drive.get('id')}")
            
            for drive in drives_data['value']:
                drive_id = drive['id']
                # If targeting a specific folder, get documents from that folder only
                if folder_path:
                    # Remove drive name from path if it's included (e.g., "Documents/Folder" -> "Folder")
                    clean_folder_path = folder_path
                    for current_drive in drives_data['value']:
                        if folder_path.startswith(f"{current_drive['name']}/"):
                            clean_folder_path = folder_path[len(f"{current_drive['name']}/"):]
                            break
                    
                    logging.info(f"📁 Targeting specific folder: {clean_folder_path} (recursive: {recursive})")
                    documents = self._get_drive_documents(drive_id, clean_folder_path, site_name, checkpoint, recursive)
                else:
                    # Start from root if no folder specified
                    documents.extend(self._get_drive_documents(drive_id, None, site_name, checkpoint, recursive))
            
            return documents
            
        except Exception as e:
            logging.error(f'Error getting documents from site {site_id}: {str(e)}')
            return []

    def _get_drive_documents(self, drive_id: str, folder_path: str = None, site_name: str = None, checkpoint: Dict[str, Any] = None, recursive: bool = False) -> List[Dict[str, Any]]:
        """Get documents from a specific drive with sequential pagination and checkpoint support"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            all_documents = []
            api_calls_made = 0
            max_api_calls = 500  # Increased from 100 to allow deeper scanning
            
            # Determine starting URL based on folder_path and checkpoint
            if folder_path:
                # If folder_path is specified, start from that specific folder
                start_url = self._build_folder_url(drive_id, folder_path)
                logging.info(f'📁 Targeting specific folder: {folder_path}')
            else:
                # Start from root if no folder specified
                start_url = f'https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children?$top=100'
                logging.info(f'📁 Starting from root folder')
            
            # Resume from checkpoint if available
            current_url = checkpoint.get('last_processed_url') if checkpoint else start_url
            if checkpoint and checkpoint.get('last_processed_url'):
                logging.info(f'📍 Resuming from checkpoint URL')
                current_url = checkpoint['last_processed_url']
            else:
                current_url = start_url
            
            # Sequential pagination through documents
            while current_url and api_calls_made < max_api_calls:
                try:
                    logging.info(f'📡 Making Graph API call #{api_calls_made + 1} - Found {len(all_documents)} documents so far')
                    response = requests.get(current_url, headers=headers)
                    response.raise_for_status()
                    api_calls_made += 1
                    
                    items = response.json()
                    items_list = items.get('value', [])
                    
                    # Process items sequentially (no randomization)
                    for item in items_list:
                        item_name = item['name']
                        item_path = self._build_item_path(folder_path or '/', item_name)
                        
                        # Check if it's a file
                        if 'file' in item:
                            file_ext = os.path.splitext(item_name)[1].lower()
                            
                            # Check if supported file type
                            if file_ext in self.supported_extensions:
                                all_documents.append({
                                    'id': item['id'],
                                    'name': item_name,
                                    'path': item_path,
                                    'download_url': item.get('@microsoft.graph.downloadUrl', ''),
                                    'last_modified': item.get('lastModifiedDateTime', ''),
                                    'size': item.get('size', 0),
                                    'folder_path': folder_path or '/',
                                    'extension': file_ext
                                })
                                logging.info(f'Found document: {item_path}')
                                
                                # Check if we've hit our discovery limit (higher than processing limit)
                                discovery_limit = self.max_documents_per_run * 10  # Find 10x more documents for better selection
                                if len(all_documents) >= discovery_limit:
                                    logging.info(f'🛑 Hit discovery limit ({discovery_limit}), will prioritize and select best documents')
                                    # Save checkpoint with next URL
                                    next_url = items.get('@odata.nextLink')
                                    if next_url and site_name:
                                        self._save_checkpoint(site_name, folder_path, next_url)
                                    return all_documents
                        
                        # Handle folder processing based on recursive flag
                        elif 'folder' in item and recursive:
                            # If recursive is True, process subfolders
                            subfolder_path = self._build_item_path(folder_path or '/', item_name)
                            logging.info(f'🔄 Processing subfolder recursively: {subfolder_path}')
                            
                            try:
                                # Get documents from subfolder
                                subfolder_documents = self._get_drive_documents(drive_id, subfolder_path, site_name, None, recursive)
                                all_documents.extend(subfolder_documents)
                                
                                # Check discovery limit after adding subfolder documents
                                discovery_limit = self.max_documents_per_run * 10
                                if len(all_documents) >= discovery_limit:
                                    logging.info(f'🛑 Hit discovery limit ({discovery_limit}) after processing subfolder')
                                    return all_documents
                                    
                            except Exception as e:
                                logging.warning(f'Error processing subfolder {subfolder_path}: {str(e)}')
                                continue
                        
                        elif 'folder' in item and not recursive:
                            # If recursive is False, skip subfolders
                            logging.info(f'⏭️ Skipping subfolder (non-recursive): {item_name}')
                    
                    # Get next page URL
                    next_url = items.get('@odata.nextLink')
                    
                    # Save checkpoint after each successful page
                    if site_name:
                        self._save_checkpoint(site_name, folder_path, next_url)
                    
                    current_url = next_url
                    
                except requests.exceptions.RequestException as e:
                    logging.error(f'Error during pagination: {str(e)}')
                    api_calls_made += 1  # Count failed calls too
                    break
            
            # Clear checkpoint if we've completed all pages
            if not current_url and site_name:
                self._save_checkpoint(site_name, folder_path, None)
                logging.info(f'✅ Completed processing all documents in folder')
            
            # Log final statistics
            logging.info(f'📊 Sequential scan completed:')
            logging.info(f'   Documents found: {len(all_documents)}')
            logging.info(f'   API calls made: {api_calls_made}')
            logging.info(f'   📋 Documents processed sequentially (no randomization)')
            
            if api_calls_made >= max_api_calls:
                logging.warning(f'⚠️ Stopped due to API call limit. Will resume from checkpoint on next run.')
            
            return all_documents
            
        except Exception as e:
            logging.error(f'Error getting documents from drive {drive_id}: {str(e)}')
            return []

    def _build_folder_url(self, drive_id: str, folder_path: str) -> str:
        """Build Graph API URL for a specific folder path"""
        # Clean up folder path
        clean_path = folder_path.strip('/')
        # Properly encode each segment
        encoded_path = '%2F'.join([quote(seg, safe='') for seg in clean_path.split('/')])
        return f'https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded_path}:/children?$top=100'
    
    def _build_item_path(self, folder_path: str, item_name: str) -> str:
        """Build full item path from folder path and item name"""
        if folder_path == '/':
            return f'/{item_name}'
        else:
            return f"{folder_path.rstrip('/')}/{item_name}"

    def _prioritize_documents_for_testing(self, documents: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Smart document prioritization with randomization to avoid processing same documents"""
        
        # *** NEW: Filter out already processed documents first ***
        unprocessed_documents = []
        already_processed_count = 0
        for doc in documents:
            if not self._is_document_processed(doc['id']):
                unprocessed_documents.append(doc)
            else:
                already_processed_count += 1
        
        logging.info(f'📊 Document status: {len(unprocessed_documents)} new documents, {already_processed_count} already processed (total: {len(documents)})')
        
        # If we have enough unprocessed documents, use those
        if len(unprocessed_documents) >= self.max_documents_per_run:
            documents_to_prioritize = unprocessed_documents
        else:
            # Use all documents if we don't have enough unprocessed ones
            documents_to_prioritize = documents
            logging.info(f'⚠️ Using all documents including processed ones (need {self.max_documents_per_run}, found {len(unprocessed_documents)} unprocessed)')
        
        def document_priority(doc):
            extension = doc.get('extension', '').lower()
            size = doc.get('size', 0)
            last_modified = doc.get('last_modified', '')
            
            # Priority score (lower = higher priority)
            priority = 0
            
            # *** NEW: Add randomization to priority ***
            priority += random.randint(0, 2)  # Add random factor to mix things up
            
            # Prioritize lightweight file types
            if extension in self.lightweight_extensions:
                priority += 1
            elif extension in ['.pdf']:
                priority += 3
            elif extension in ['.xlsx', '.xls']:
                priority += 4
            else:
                priority += 2
            
            # Prioritize smaller files
            if size < 1024 * 1024:  # < 1MB
                priority += 1
            elif size > 10 * 1024 * 1024:  # > 10MB
                priority += 5
            
            # Prioritize recently modified (parse date if available)
            try:
                if last_modified:
                    mod_date = datetime.datetime.fromisoformat(last_modified.replace('Z', '+00:00'))
                    days_old = (datetime.datetime.utcnow().replace(tzinfo=mod_date.tzinfo) - mod_date).days
                    if days_old < 30:  # Modified in last 30 days
                        priority += 1
            except:
                pass
            
            # *** NEW: Boost priority for unprocessed documents ***
            if not self._is_document_processed(doc['id']):
                priority -= 5  # Higher priority for unprocessed docs
            
            return priority
        
        # Sort by priority and take the first max_documents_per_run
        sorted_docs = sorted(documents_to_prioritize, key=document_priority)
        selected_docs = sorted_docs[:self.max_documents_per_run]
        
        # *** NEW: Final randomization within selected documents ***
        random.shuffle(selected_docs)
        
        logging.info(f'🎲 Selected {len(selected_docs)} documents with randomization and priority')
        
        return selected_docs

    def _is_document_processed(self, doc_id: str) -> bool:
        """Check if document has already been processed by looking for first chunk"""
        try:
            blob_client = self.storage_client.get_blob_client(
                container="jennifur-processed",
                blob=f"{doc_id}_0.json"
            )
            return blob_client.exists()
        except Exception:
            return False

    def _download_document(self, download_url: str) -> bytes:
        """Download document content from OneDrive/SharePoint"""
        try:
            headers = {'Authorization': f'Bearer {self.graph_token}'}
            response = requests.get(download_url, headers=headers)
            response.raise_for_status()
            return response.content
            
        except Exception as e:
            logging.error(f'Error downloading document: {str(e)}')
            raise

    def _extract_text_with_cost_tracking(self, doc_content: bytes, filename: str) -> tuple[str, float, bool]:
        """Extract text and track estimated costs. Returns (text, cost, success_flag)"""
        
        # Validate document first
        if not self._validate_document_before_processing(doc_content, filename):
            return self._create_fallback_text_content(filename, "File validation failed"), 0.0, False
        
        file_ext = os.path.splitext(filename)[1].lower()
        file_size_mb = len(doc_content) / (1024 * 1024)
        
        # Estimated cost calculation (based on Azure Document Intelligence pricing)
        # Approximately $1.00 per 1,000 pages, estimate ~0.5MB per page
        estimated_pages = max(1, file_size_mb / 0.5)
        estimated_cost = estimated_pages * 0.001  # $0.001 per page
        
        try:
            logging.info(f'Analyzing document with Document Intelligence: {filename} (Est. cost: ${estimated_cost:.4f})')
            
            # Analyze document
            poller = self.doc_intelligence_client.begin_analyze_document(
                "prebuilt-document", 
                doc_content
            )
            result = poller.result()
            
            # Extract content based on file type
            if file_ext in ['.xlsx', '.xls']:
                extracted_text = self._extract_excel_content(result, filename)
            else:
                extracted_text = self._extract_standard_content(result, filename)
            
            return extracted_text, estimated_cost, True
            
        except Exception as e:
            error_message = str(e)
            
            # PowerPoint fallback: try extracting text locally if DI fails
            if file_ext in ['.pptx', '.ppt']:
                try:
                    prs = Presentation(BytesIO(doc_content))
                    pptx_text = f"PowerPoint Document: {filename}\n"
                    pptx_text += f"Processed with: python-pptx fallback (DI unavailable)\n"
                    pptx_text += f"Processing date: {datetime.datetime.utcnow().isoformat()}\n\n"
                    for i, slide in enumerate(prs.slides):
                        pptx_text += f"--- Slide {i+1} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                pptx_text += shape.text + "\n"
                        pptx_text += "\n"
                    if pptx_text.strip():
                        logging.info(f'Successfully extracted {len(pptx_text):,} characters from PowerPoint file {filename} using python-pptx fallback')
                        return pptx_text, 0.0, True
                except Exception as pptx_e:
                    logging.warning(f'python-pptx fallback failed for {filename}: {pptx_e}')
                    # Fall through to existing fallback logic
            
            # Excel fallback: try extracting text locally if DI fails and file is .xlsx
            if file_ext == '.xlsx':
                try:
                    wb = openpyxl.load_workbook(BytesIO(doc_content), data_only=True)
                    excel_text = f"Excel Document: {filename}\n"
                    excel_text += f"Processed with: openpyxl fallback (DI unavailable)\n"
                    excel_text += f"Processing date: {datetime.datetime.utcnow().isoformat()}\n\n"
                    
                    for sheet in wb.worksheets:
                        sheet_has_content = False
                        sheet_content = f"--- Sheet: {sheet.title} ---\n"
                        
                        # Only process rows with actual data (skip empty rows)
                        for row in sheet.iter_rows(values_only=True):
                            # Filter out None/empty cells and join with pipe separators
                            row_data = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                            if row_data:  # Only add non-empty rows
                                if not sheet_has_content:
                                    excel_text += sheet_content
                                    sheet_has_content = True
                                excel_text += ' | '.join(row_data) + ' | \n'
                        
                        if sheet_has_content:
                            excel_text += '\n'
                    
                    if excel_text.strip():
                        logging.info(f'Successfully extracted {len(excel_text):,} characters from Excel file {filename} using openpyxl fallback (filtered empty cells)')
                        return excel_text, 0.0, True
                except Exception as xlsx_e:
                    logging.warning(f'openpyxl fallback failed for {filename}: {xlsx_e}')
                    # Fall through to existing fallback logic
                    
            # Handle specific Document Intelligence errors
            if any(phrase in error_message for phrase in [
                "InvalidContent", "corrupted or format is unsupported",
                "UnsupportedMediaType", "InvalidRequest"
            ]):
                logging.warning(f'⚠️ Document Intelligence cannot process {filename}: {error_message}')
                return self._create_fallback_text_content(filename, f"Document Intelligence error: {error_message}"), 0.0, False
            else:
                logging.error(f'Unexpected error extracting text from {filename}: {error_message}')
                return self._create_fallback_text_content(filename, f"Processing error: {error_message}"), 0.0, False

    def _extract_excel_content(self, result, filename: str) -> str:
        """Extract content specifically formatted for Excel files with empty cell filtering"""
        
        extracted_text = f"Excel Document: {filename}\n"
        extracted_text += f"Processed with: Azure Document Intelligence (Excel Mode)\n"
        extracted_text += f"Processing date: {datetime.datetime.utcnow().isoformat()}\n\n"
        
        # Handle Excel tables with better formatting and empty cell filtering
        if result.tables:
            for table_idx, table in enumerate(result.tables):
                sheet_content = f"--- Sheet: {table_idx + 1} ---\n"
                
                # Group cells by row and filter out empty content
                rows_dict = {}
                for cell in table.cells:
                    cell_content = cell.content.strip()
                    if cell_content:  # Only process non-empty cells
                        if cell.row_index not in rows_dict:
                            rows_dict[cell.row_index] = {}
                        rows_dict[cell.row_index][cell.column_index] = cell_content
                
                # Only add sheet if it has content
                if rows_dict:
                    extracted_text += sheet_content
                    
                    # Process rows that have content
                    for row_idx in sorted(rows_dict.keys()):
                        row_cells = rows_dict[row_idx]
                        if row_cells:  # Only process non-empty rows
                            # Create row with only non-empty cells
                            row_values = []
                            for col_idx in sorted(row_cells.keys()):
                                row_values.append(row_cells[col_idx])
                            
                            if row_values:  # Only add if row has content
                                extracted_text += " | ".join(row_values) + " | \n"
                    
                    extracted_text += "\n"
        
        # Also extract any general text content
        text_content = ""
        for page in result.pages:
            for line in page.lines:
                line_content = line.content.strip()
                if line_content:  # Only add non-empty lines
                    text_content += line_content + "\n"
        
        if text_content.strip():
            extracted_text += "--- Additional Text Content ---\n"
            extracted_text += text_content + "\n"
        
        logging.info(f'Successfully extracted {len(extracted_text):,} characters from Excel file {filename} (filtered empty cells)')
        return extracted_text

    def _extract_standard_content(self, result, filename: str) -> str:
        """Extract content for standard documents (PDF, Word, etc.)"""
        
        extracted_text = f"Document: {filename}\n"
        extracted_text += f"Processed with: Azure Document Intelligence\n"
        extracted_text += f"Processing date: {datetime.datetime.utcnow().isoformat()}\n\n"
        
        # Add page content
        for page_idx, page in enumerate(result.pages):
            extracted_text += f"--- Page {page_idx + 1} ---\n"
            for line in page.lines:
                extracted_text += line.content + "\n"
            extracted_text += "\n"
        
        # Add tables if present
        if result.tables:
            extracted_text += "--- Tables ---\n"
            for table_idx, table in enumerate(result.tables):
                extracted_text += f"Table {table_idx + 1}:\n"
                for cell in table.cells:
                    extracted_text += f"Row {cell.row_index}, Col {cell.column_index}: {cell.content}\n"
                extracted_text += "\n"
        
        logging.info(f'Successfully extracted {len(extracted_text):,} characters from {filename}')
        return extracted_text

    def _validate_document_before_processing(self, doc_content: bytes, filename: str) -> bool:
        """Validate document before sending to Document Intelligence"""
        
        # Check file size (Document Intelligence has limits)
        max_size = 500 * 1024 * 1024  # 500MB limit
        if len(doc_content) > max_size:
            logging.warning(f'File {filename} too large ({len(doc_content):,} bytes), skipping Document Intelligence')
            return False
        
        # Check for minimum file size
        if len(doc_content) < 100:
            logging.warning(f'File {filename} too small ({len(doc_content)} bytes), might be corrupted')
            return False
        
        # Check for empty files
        if len(doc_content) == 0:
            logging.warning(f'File {filename} is empty, skipping')
            return False
        
        return True

    def _create_fallback_text_content(self, filename: str, reason: str) -> str:
        """Create fallback text content when Document Intelligence fails"""
        
        fallback_text = f"Document: {filename}\n"
        fallback_text += f"Processed with: Fallback (Document Intelligence unavailable)\n"
        fallback_text += f"Processing date: {datetime.datetime.utcnow().isoformat()}\n"
        fallback_text += f"Reason: {reason}\n\n"
        fallback_text += f"--- Document Information ---\n"
        fallback_text += f"Filename: {filename}\n"
        fallback_text += f"Note: This document could not be fully processed. "
        fallback_text += f"The document exists and was detected, but text extraction was not possible.\n"
        fallback_text += f"File may contain images, complex formatting, or be in an unsupported format.\n"
        
        return fallback_text

    def _extract_client_metadata_from_path(self, doc_path: str) -> Dict[str, Any]:
        """Extract client and PM metadata from SharePoint folder path, prioritizing main client folder"""
        import re
        
        logging.info(f'🏢 Extracting client metadata from path: {doc_path}')
        
        # First, look for the main client folder pattern: Client Name (PM-X)/
        # This ensures we extract from the main folder, not subfolders like "Archived"
        client_folder_pattern = r'([^/]+)\s*\(PM-([A-Za-z])\)/'
        match = re.search(client_folder_pattern, doc_path, re.IGNORECASE)
        
        if match:
            client_name = match.group(1).strip()
            pm_code = match.group(2).upper()
            
            logging.info(f'✅ Found main client folder: "{client_name}" with PM-{pm_code}')
        else:
            # Fallback: look for client pattern anywhere in path (less preferred)
            fallback_patterns = [
                r'([^/]+)\s*\(PM-([A-Z])\)',  # Standard format
                r'([^/]+)\s*\(PM-([a-zA-Z])\)',  # Case insensitive PM code
                r'(.*?)\s*\(PM-([A-Z])\)',  # Any text before PM pattern
            ]
            
            for pattern in fallback_patterns:
                match = re.search(pattern, doc_path, re.IGNORECASE)
                if match:
                    client_name = match.group(1).strip()
                    pm_code = match.group(2).upper()
                    logging.info(f'⚠️ Using fallback pattern for client: "{client_name}" with PM-{pm_code}')
                    break
            else:
                logging.warning(f'❌ No client metadata found in path: {doc_path}')
                return None
        
        # Map PM codes to names
        pm_mapping = {
            'C': 'Caleb',
            'K': 'Katherine', 
            'S': 'Sam'
        }
        
        pm_name = pm_mapping.get(pm_code, pm_code)
        
        # Clean up client name (remove leading/trailing slashes and whitespace)
        client_name = client_name.strip('/ ')
        
        result = {
            'client_name': client_name,
            'pm_code': pm_code,
            'pm_name': pm_name,
            'confidence_score': 0.9
        }
        
        logging.info(f'🎯 Extracted client metadata: {result}')
        return result

    def _process_magic_meeting_tracker(self, doc_content: bytes, doc_name: str, doc_path: str, doc_id: str) -> List[Dict[str, Any]]:
        """Special processing for Magic Meeting Tracker Excel file with sheet-based client attribution"""
        try:
            # Import the enhanced Excel processor
            import sys
            import os
            sys.path.append(os.path.dirname(os.path.dirname(__file__)))
            from enhanced_excel_processor import EnhancedExcelProcessor
            
            logging.info(f'🎯 Processing Magic Meeting Tracker with sheet-based client detection')
            
            # Initialize the enhanced processor
            processor = EnhancedExcelProcessor()
            excel_data = processor.extract_from_excel(doc_content, doc_name)
            
            if excel_data.get("type") == "excel_error":
                logging.error(f'Enhanced Excel processing failed: {excel_data.get("error")}')
                # Fallback to standard processing
                return self._fallback_excel_processing(doc_content, doc_name, doc_path, doc_id)
            
            # Create chunks for each sheet with its own client metadata
            all_chunks = []
            current_timestamp = datetime.datetime.utcnow().isoformat() + "Z"
            
            for sheet_name, sheet_data in excel_data["sheets"].items():
                content = sheet_data["content"]
                client_info = sheet_data.get("client_info")
                
                if not content or len(content.strip()) < 10:  # Lower minimum content threshold
                    logging.warning(f'⚠️ Sheet "{sheet_name}" skipped due to insufficient content: {len(content) if content else 0} characters')
                    continue
                
                logging.info(f'📄 Processing sheet "{sheet_name}": {len(content)} characters, client_info: {bool(client_info)}')
                
                # Determine client metadata for this sheet
                if client_info and client_info.get("sheet_client_name"):
                    # Use detected client from sheet name
                    sheet_client_name = client_info["sheet_client_name"]
                    sheet_pm_initial = client_info.get("sheet_pm_initial", "N/A")
                    sheet_pm_name = client_info.get("sheet_pm_name", "N/A")
                    is_client_specific = True
                    has_client_folder = True
                    
                    logging.info(f'📋 Sheet "{sheet_name}" attributed to client: {sheet_client_name}')
                else:
                    # Default to internal if no client detected from sheet name
                    sheet_client_name = "Autobahn Internal"
                    sheet_pm_initial = "N/A"
                    sheet_pm_name = "N/A"
                    is_client_specific = False
                    has_client_folder = False
                    
                    logging.warning(f'⚠️ Sheet "{sheet_name}" could not detect client, defaulting to Autobahn Internal')
                
                # Create metadata for this sheet
                sheet_metadata = {
                    "document_path": f"{doc_path} (Sheet: {sheet_name})",
                    "filename": f"{doc_name} - {sheet_name}",
                    "processed_timestamp": current_timestamp,
                    "content_length": len(content),
                    "word_count": len(content.split()),
                    "character_count": len(content),
                    "folder_depth": len([p for p in doc_path.split('/') if p.strip()]),
                    "file_extension": ".xlsx",
                    "is_client_specific": is_client_specific,
                    "has_client_folder": has_client_folder,
                    "client_name": sheet_client_name,
                    "pm_initial": sheet_pm_initial,
                    "pm_name": sheet_pm_name,
                    "document_category": "meeting_tracking",
                    "document_id": f"{doc_id}_sheet_{sheet_name.replace(' ', '_')}",
                    "chunk_method": "enhanced_excel_sheet_based",
                    "source": "Magic Meeting Tracker",
                    "sheet_name": sheet_name,
                    "sheet_client_confidence": client_info.get("confidence", 0.0) if client_info else 0.0
                }
                
                # Chunk the sheet content (will be mostly single chunk unless very large)
                sheet_chunks = self._chunk_excel_sheet_content(content, sheet_metadata)
                all_chunks.extend(sheet_chunks)
            
            logging.info(f'✅ Successfully processed Magic Meeting Tracker: {len(all_chunks)} chunks from {len(excel_data["sheets"])} sheets')
            return all_chunks
            
        except Exception as e:
            logging.error(f'Error processing Magic Meeting Tracker: {str(e)}')
            # Fallback to standard processing
            return self._fallback_excel_processing(doc_content, doc_name, doc_path, doc_id)
    
    def _chunk_excel_sheet_content(self, content: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Chunk content from a single Excel sheet"""
        chunk_size = 1000
        overlap = 100
        chunks = []
        
        if len(content) <= chunk_size:
            # Single chunk
            chunk_data = {
                **metadata,
                "chunk": content,
                "chunk_id": f"{metadata['document_id']}_0",
                "parent_id": metadata['document_id'],
                "chunk_index": 0,
                "title": f"{metadata['filename']} - {metadata['sheet_name']}"
            }
            chunks.append(chunk_data)
        else:
            # Multiple chunks with overlap
            start = 0
            chunk_index = 0
            
            while start < len(content):
                end = min(start + chunk_size, len(content))
                chunk_content = content[start:end]
                
                chunk_data = {
                    **metadata,
                    "chunk": chunk_content,
                    "chunk_id": f"{metadata['document_id']}_{chunk_index}",
                    "parent_id": metadata['document_id'],
                    "chunk_index": chunk_index,
                    "title": f"{metadata['filename']} - {metadata['sheet_name']} (Part {chunk_index + 1})"
                }
                chunks.append(chunk_data)
                
                # Move start position with overlap
                start += chunk_size - overlap
                chunk_index += 1
        
        return chunks
    
    def _fallback_excel_processing(self, doc_content: bytes, doc_name: str, doc_path: str, doc_id: str) -> List[Dict[str, Any]]:
        """Fallback Excel processing using standard method"""
        try:
            extracted_text, cost, success = self._extract_text_with_cost_tracking(doc_content, doc_name)
            if success and extracted_text:
                client_metadata = self._extract_client_metadata_from_path(doc_path)
                chunks = self._chunk_text(extracted_text, doc_id, doc_name, doc_path, client_metadata)
                return chunks
        except Exception as e:
            logging.error(f'Fallback Excel processing failed: {str(e)}')
        return []
    
    def _chunk_text(self, text: str, doc_id: str, doc_name: str, doc_path: str, client_metadata: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """Chunk text into 1000 character chunks with 100 character overlap, matching specified format"""
        
        chunk_size = 1000
        overlap = 100
        chunks = []
        
        # Calculate document statistics
        word_count = len(text.split())
        character_count = len(text)
        
        # Extract folder depth from path
        folder_depth = len([p for p in doc_path.split('/') if p.strip()]) - 1
        
        # Get file extension
        file_extension = os.path.splitext(doc_name)[1].lower()
        
        # Determine if this is client specific and categorize
        is_client_specific = bool(client_metadata)
        has_client_folder = is_client_specific
        
        # Determine document category based on filename/path
        document_category = self._categorize_document(doc_name, doc_path)
        
        # Set client info
        if client_metadata:
            client_name = client_metadata.get('client_name', 'Unknown')
            pm_initial = client_metadata.get('pm_code', 'N/A')
            pm_name = client_metadata.get('pm_name', 'N/A')
        else:
            client_name = "Autobahn Internal"
            pm_initial = "N/A"
            pm_name = "N/A"
        
        current_timestamp = datetime.datetime.utcnow().isoformat() + "Z"
        
        # If text is shorter than chunk size, return single chunk
        if len(text) <= chunk_size:
            chunk_data = {
                "document_path": doc_path,
                "filename": doc_name,
                "processed_timestamp": current_timestamp,
                "content_length": character_count,
                "client_name": client_name,
                "pm_initial": pm_initial,
                "pm_name": pm_name,
                "document_category": document_category,
                "is_client_specific": is_client_specific,
                "metadata_updated_timestamp": current_timestamp,
                "folder_depth": folder_depth,
                "file_extension": file_extension,
                "has_client_folder": has_client_folder,
                "word_count": word_count,
                "character_count": character_count,
                "chunk": text,
                "chunk_id": f"{doc_id}_0",
                "chunk_index": 0,
                "parent_id": f"{doc_id}",
                "processing_method": "reprocessed_enhanced"
            }
            chunks.append(chunk_data)
            return chunks
        
        # Create overlapping chunks
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunk_content = text[start:end]
            
            # Try to break at word boundary if we're not at the end
            if end < len(text):
                last_space = chunk_content.rfind(' ')
                if last_space > chunk_size * 0.8:  # Only break if we're not losing too much content
                    end = start + last_space
                    chunk_content = text[start:end]
            
            chunk_data = {
                "document_path": doc_path,
                "filename": doc_name,
                "processed_timestamp": current_timestamp,
                "content_length": character_count,
                "client_name": client_name,
                "pm_initial": pm_initial,
                "pm_name": pm_name,
                "document_category": document_category,
                "is_client_specific": is_client_specific,
                "metadata_updated_timestamp": current_timestamp,
                "folder_depth": folder_depth,
                "file_extension": file_extension,
                "has_client_folder": has_client_folder,
                "word_count": word_count,
                "character_count": character_count,
                "chunk": chunk_content,
                "chunk_id": f"{doc_id}_{chunk_index}",
                "chunk_index": chunk_index,
                "parent_id": f"{doc_id}",
                "processing_method": "reprocessed_enhanced"
            }
            
            chunks.append(chunk_data)
            
            # Move to next chunk with overlap
            start = max(start + chunk_size - overlap, end)
            chunk_index += 1
            
            # Safety check to prevent infinite loops
            if chunk_index > 1000:  # Max 1000 chunks per document
                logging.warning(f"Hit maximum chunk limit for document {doc_name}")
                break
        
        logging.info(f"Created {len(chunks)} chunks for document {doc_name}")
        return chunks

    def _categorize_document(self, filename: str, doc_path: str) -> str:
        """Categorize document based on path structure and filename"""
        
        # First try to extract from path structure (second level after client name)
        path_category = self._extract_document_category_from_path(doc_path)
        if path_category and path_category != "general":
            return path_category
        
        # Fallback to filename-based categorization
        filename_lower = filename.lower()
        path_lower = doc_path.lower()
        
        # Financial documents
        if any(term in filename_lower or term in path_lower for term in [
            'financial', 'finance', 'budget', 'cost', 'invoice', 'payment', 
            'revenue', 'profit', 'loss', 'expense', 'accounting', 'packet'
        ]):
            return "financials"
        
        # Roadmap documents  
        if any(term in filename_lower or term in path_lower for term in [
            'roadmap', 'timeline', 'schedule', 'plan', 'milestone', 'strategy'
        ]):
            return "roadmaps"
        
        # Organization charts
        if any(term in filename_lower or term in path_lower for term in [
            'org', 'organization', 'orgchart', 'org chart', 'hierarchy', 'structure'
        ]):
            return "org_charts"
        
        # Contracts and legal
        if any(term in filename_lower or term in path_lower for term in [
            'contract', 'agreement', 'legal', 'terms', 'sow', 'statement of work'
        ]):
            return "contracts"
        
        # Reports
        if any(term in filename_lower or term in path_lower for term in [
            'report', 'analysis', 'summary', 'status', 'update'
        ]):
            return "reports"
        
        # Templates
        if any(term in filename_lower or term in path_lower for term in [
            'template', 'example', 'sample', 'format'
        ]):
            return "templates"
        
        # Default category
        return "general"

    def _extract_document_category_from_path(self, doc_path: str) -> str:
        """Extract document category from the second level of the path after client name"""
        import re
        
        # First, find the client folder pattern
        client_pattern = r'([^/]+)\s*\(PM-[A-Z]\)/([^/]+)'
        match = re.search(client_pattern, doc_path, re.IGNORECASE)
        
        if match:
            category_folder = match.group(2).strip()
            
            # Extract category from patterns like "04. Roadmaps & Org Charts" -> "roadmaps_org_charts"
            # Pattern: number followed by dot and words
            category_pattern = r'^\d+\.\s*(.+)'
            category_match = re.search(category_pattern, category_folder, re.IGNORECASE)
            
            if category_match:
                category_text = category_match.group(1).strip()
                # Clean and normalize the category
                # Remove special characters, convert to lowercase, replace spaces with underscores
                clean_category = re.sub(r'[^\w\s]', '', category_text)  # Remove special chars
                clean_category = re.sub(r'\s+', '_', clean_category.strip().lower())  # Spaces to underscores, lowercase
                
                return clean_category
            else:
                # If no number pattern, use the folder name directly (cleaned)
                clean_category = re.sub(r'[^\w\s]', '', category_folder)
                clean_category = re.sub(r'\s+', '_', clean_category.strip().lower())
                
                if clean_category:
                    return clean_category
        
        return "general"

    def _store_processed_document_with_chunks(self, doc_id: str, filename: str, content: str, doc_metadata: Dict[str, Any], doc_path: str, chunks: List[Dict[str, Any]]) -> None:
        """Store individual chunks as separate blobs in jennifur-processed container"""
        try:
            # Store each chunk as a separate blob
            for chunk in chunks:
                chunk_id = chunk["chunk_id"]
                blob_name = f"{chunk_id}.json"
                
                blob_client = self.storage_client.get_blob_client(
                    container="jennifur-processed",
                    blob=blob_name
                )
                
                # Upload the chunk data
                blob_client.upload_blob(
                    json.dumps(chunk, indent=2),
                    overwrite=True,
                    content_type='application/json'
                )
            
            logging.info(f'Document {filename} stored in jennifur-processed container with {len(chunks)} individual chunk files')
            
        except Exception as e:
            logging.error(f'Error storing processed document: {str(e)}')
            raise
    
    def _store_magic_meeting_tracker_chunks(self, chunks: List[Dict[str, Any]], doc_id: str, filename: str, doc_path: str, doc: Dict[str, Any]):
        """Store Magic Meeting Tracker chunks with proper client attribution"""
        try:
            container_client = self.storage_client.get_container_client("jennifur-processed")
            
            for i, chunk in enumerate(chunks):
                chunk_filename = f"{chunk['chunk_id']}.json"
                
                # Add common document metadata
                chunk.update({
                    "original_filename": filename,
                    "original_document_id": doc_id,
                    "download_url": doc.get('download_url', ''),
                    "file_size_bytes": doc.get('size', 0),
                    "sharepoint_item_id": doc.get('id', ''),
                    "last_modified": doc.get('last_modified', ''),
                    "processing_method": "magic_meeting_tracker_specialized"
                })
                
                # Upload chunk
                container_client.upload_blob(
                    name=chunk_filename,
                    data=json.dumps(chunk, indent=2),
                    overwrite=True,
                    content_type='application/json'
                )
                
                logging.info(f'📄 Stored chunk {i+1}: {chunk_filename} (client: {chunk.get("client_name", "Unknown")}, sheet: {chunk.get("sheet_name", "Unknown")})')
            
            logging.info(f'✅ Magic Meeting Tracker stored: {len(chunks)} chunks from {filename}')
            
        except Exception as e:
            logging.error(f'Error storing Magic Meeting Tracker chunks: {str(e)}')
            raise

    def _process_single_document_with_cost_control(self, doc: Dict[str, Any]) -> Dict[str, Any]:
        """Process a single document with enhanced cost control and size limits"""
        
        doc_id = doc['id']
        doc_name = doc['name']
        doc_path = doc.get('path', doc_name)
        doc_extension = doc.get('extension', 'unknown')
        doc_size = doc.get('size', 0)
        
        processing_start = datetime.datetime.utcnow()
        estimated_cost = 0.0
        
        try:
            # Check if already processed
            if self._is_document_processed(doc_id):
                logging.info(f'📄✅ Document already processed, skipping: {doc_path} (found existing chunk file {doc_id}_0.json)')
                return {
                    "action": "skipped", 
                    "reason": "already_processed", 
                    "path": doc_path,
                    "extension": doc_extension,
                    "cost_estimate": 0.0
                }
            
            # Size-based processing limits
            max_size_bytes = self.max_file_size_mb * 1024 * 1024
            if doc_size > max_size_bytes:
                logging.info(f'📏 File {doc_path} exceeds size limit ({doc_size:,} bytes), skipping')
                return {
                    "action": "skipped_size_limit",
                    "reason": f"file_too_large_{doc_size}_bytes",
                    "path": doc_path,
                    "extension": doc_extension,
                    "cost_estimate": 0.0
                }
            
            logging.info(f'Processing document: {doc_path} ({doc_extension}, {doc_size:,} bytes)')
            
            # Download document with error handling
            try:
                doc_content = self._download_document(doc['download_url'])
                download_size = len(doc_content)
                logging.info(f'Downloaded {download_size:,} bytes for {doc_name}')
                
            except Exception as e:
                logging.error(f'Failed to download {doc_path}: {str(e)}')
                return {
                    "action": "error", 
                    "reason": f"download_failed: {str(e)}", 
                    "path": doc_path,
                    "extension": doc_extension,
                    "cost_estimate": 0.0
                }
            
            # Special handling for Magic Meeting Tracker Excel file
            if doc_name.lower().startswith('magic meeting tracker') and doc_extension in ['.xlsx', '.xls']:
                logging.info(f'🎯 Detected Magic Meeting Tracker, using specialized processing')
                try:
                    chunks = self._process_magic_meeting_tracker(doc_content, doc_name, doc_path, doc_id)
                    if chunks:
                        # Store chunks directly
                        self._store_magic_meeting_tracker_chunks(chunks, doc_id, doc_name, doc_path, doc)
                        
                        result_base = {
                            "path": doc_path,
                            "extension": doc_extension,
                            "processing_duration_seconds": round((datetime.datetime.utcnow() - processing_start).total_seconds(), 2),
                            "file_size_bytes": download_size,
                            "cost_estimate": round(estimated_cost, 4),
                            "chunk_count": len(chunks),
                            "sheets_processed": len(set(chunk.get('sheet_name', '') for chunk in chunks))
                        }
                        logging.info(f'✅ Magic Meeting Tracker processed: {len(chunks)} chunks from {result_base["sheets_processed"]} sheets')
                        return {**result_base, "action": "processed"}
                    else:
                        logging.warning(f'⚠️ Magic Meeting Tracker processing returned no chunks, falling back to standard processing')
                except Exception as e:
                    logging.error(f'Magic Meeting Tracker processing failed: {str(e)}, falling back to standard processing')
            
            # Extract text with cost estimation
            try:
                extracted_text, text_extraction_cost, extraction_success = self._extract_text_with_cost_tracking(doc_content, doc_name)
                estimated_cost += text_extraction_cost
                
                # Add folder path context and client metadata to extracted text
                folder_context = f"Document Location: {doc_path}\n"
                folder_context += f"Folder: {doc.get('folder_path', '/')}\n"
                folder_context += f"File Extension: {doc_extension}\n"
                folder_context += f"File Size: {download_size:,} bytes\n"
                
                # Extract client metadata from folder path
                client_metadata = self._extract_client_metadata_from_path(doc_path)
                if client_metadata:
                    folder_context += f"Client: {client_metadata.get('client_name', 'Unknown')}\n"
                    folder_context += f"PM: {client_metadata.get('pm_name', 'Unknown')} (PM-{client_metadata.get('pm_code', 'N/A')})\n"
                    logging.info(f'📋 Using client metadata: {client_metadata}')
                else:
                    logging.warning(f'⚠️ No client metadata found for document: {doc_path}')
                
                folder_context += "\n"
                extracted_text = folder_context + extracted_text
                
            except Exception as e:
                logging.error(f'Text extraction completely failed for {doc_path}: {str(e)}')
                extracted_text = self._create_fallback_text_content(doc_name, f"Complete extraction failure: {str(e)}")
                extraction_success = False
            
            # Handle failed extractions - skip instead of processing
            if not extraction_success:
                logging.info(f'⏭️ Skipping document {doc_path} - Document Intelligence extraction failed')
                result_base = {
                    "path": doc_path,
                    "extension": doc_extension,
                    "processing_duration_seconds": round((datetime.datetime.utcnow() - processing_start).total_seconds(), 2),
                    "file_size_bytes": download_size if 'download_size' in locals() else 0,
                    "cost_estimate": round(estimated_cost, 4)
                }
                return {**result_base, "action": "skipped_extraction_failed", "reason": "document_intelligence_failed"}
            
            # Chunk the document text
            chunks = self._chunk_text(extracted_text, doc_id, doc_name, doc_path, client_metadata)
            
            # Store the processed document with chunks
            try:
                self._store_processed_document_with_chunks(doc_id, doc_name, extracted_text, doc, doc_path, chunks)
                processing_duration = (datetime.datetime.utcnow() - processing_start).total_seconds()
                
                logging.info(f'✅ Document {doc_path} processed successfully with {len(chunks)} chunks (${estimated_cost:.4f})')
                return {
                    "action": "processed",
                    "reason": "successfully_processed_and_chunked",
                    "path": doc_path,
                    "extension": doc_extension,
                    "processing_duration_seconds": round(processing_duration, 2),
                    "file_size_bytes": download_size if 'download_size' in locals() else 0,
                    "cost_estimate": round(estimated_cost, 4),
                    "chunks_created": len(chunks),
                    "content_length": len(extracted_text)
                }
            except Exception as e:
                logging.error(f'Failed to store processed document {doc_path}: {str(e)}')
                return {
                    "action": "error",
                    "reason": f"storage_failed: {str(e)}",
                    "path": doc_path,
                    "extension": doc_extension,
                    "cost_estimate": round(estimated_cost, 4)
                }
            
        except Exception as e:
            processing_duration = (datetime.datetime.utcnow() - processing_start).total_seconds()
            logging.error(f'Unexpected error processing document {doc_path}: {str(e)}')
            return {
                "action": "error", 
                "reason": f"unexpected_error: {str(e)}", 
                "path": doc_path,
                "extension": doc_extension,
                "processing_duration_seconds": round(processing_duration, 2),
                "cost_estimate": round(estimated_cost, 4)
            }

    def process_single_file(self, site_name: str, folder_path: str, file_name: str) -> Dict[str, Any]:
        """Process a single specific file"""
        return {
            "status": "success",
            "message": f"Single file processing requested for {file_name}",
            "site_name": site_name,
            "folder_path": folder_path,
            "file_name": file_name
        }