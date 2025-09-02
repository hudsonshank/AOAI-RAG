"""
Enhanced Document Processor with Modern RAG Best Practices

Provides async processing, advanced chunking, agentic capabilities, and production-ready patterns
inspired by leading RAG implementations.
"""

import asyncio
import hashlib
import json
import logging
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple, AsyncIterator
from dataclasses import dataclass, asdict
from enum import Enum
import aiohttp
from concurrent.futures import ThreadPoolExecutor
from azure.storage.blob.aio import BlobServiceClient
from azure.ai.formrecognizer.aio import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from azure.identity.aio import DefaultAzureCredential

# Import your existing components
from ..utils.enhanced_excel_processor import EnhancedExcelProcessor
from ..utils.client_metadata_extractor import ClientMetadataExtractor


class ProcessingStatus(Enum):
    """Document processing status"""
    PENDING = "pending"
    PROCESSING = "processing" 
    COMPLETED = "completed"
    FAILED = "failed"
    SKIPPED = "skipped"


class DocumentType(Enum):
    """Supported document types"""
    PDF = "pdf"
    WORD = "word" 
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    TEXT = "text"
    UNKNOWN = "unknown"


@dataclass
class DocumentMetadata:
    """Enhanced document metadata structure"""
    id: str
    filename: str
    path: str
    size: int
    type: DocumentType
    processing_status: ProcessingStatus
    client_name: str
    pm_name: str
    document_category: str
    is_client_specific: bool
    processed_timestamp: Optional[str] = None
    error_message: Optional[str] = None
    chunk_count: Optional[int] = None
    confidence_scores: Optional[Dict[str, float]] = None


@dataclass 
class ProcessingResult:
    """Processing result with detailed metrics"""
    document_id: str
    status: ProcessingStatus
    chunks_created: int = 0
    processing_time_ms: int = 0
    tokens_processed: int = 0
    confidence_score: float = 0.0
    error_message: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None


class EnhancedDocumentProcessor:
    """
    Production-ready document processor with async capabilities,
    advanced chunking, and agentic processing patterns.
    """
    
    def __init__(self, 
                 storage_connection: str,
                 doc_intelligence_endpoint: str,
                 doc_intelligence_key: str,
                 max_concurrent_docs: int = 10,
                 chunk_size: int = 1000,
                 chunk_overlap: int = 100):
        """Initialize the enhanced processor"""
        
        self.storage_connection = storage_connection
        self.doc_intelligence_endpoint = doc_intelligence_endpoint
        self.doc_intelligence_key = doc_intelligence_key
        self.max_concurrent_docs = max_concurrent_docs
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # Initialize components
        self.logger = logging.getLogger(__name__)
        self.excel_processor = EnhancedExcelProcessor()
        self.client_extractor = ClientMetadataExtractor()
        
        # Thread pool for CPU-intensive tasks
        self.executor = ThreadPoolExecutor(max_workers=4)
        
        # Processing statistics
        self.stats = {
            "documents_processed": 0,
            "total_chunks_created": 0,
            "processing_errors": 0,
            "start_time": datetime.utcnow()
        }
        
        self.logger.info("Enhanced Document Processor initialized")
    
    async def process_documents_batch(self, 
                                    documents: List[Dict[str, Any]],
                                    batch_size: int = 10) -> AsyncIterator[ProcessingResult]:
        """
        Process documents in batches with async concurrency control
        """
        semaphore = asyncio.Semaphore(self.max_concurrent_docs)
        
        async def process_single_with_semaphore(doc):
            async with semaphore:
                return await self.process_single_document(doc)
        
        # Process in batches
        for i in range(0, len(documents), batch_size):
            batch = documents[i:i + batch_size]
            self.logger.info(f"Processing batch {i//batch_size + 1}: {len(batch)} documents")
            
            # Process batch concurrently
            tasks = [process_single_with_semaphore(doc) for doc in batch]
            results = await asyncio.gather(*tasks, return_exceptions=True)
            
            for result in results:
                if isinstance(result, Exception):
                    self.logger.error(f"Batch processing error: {result}")
                    yield ProcessingResult(
                        document_id="unknown",
                        status=ProcessingStatus.FAILED,
                        error_message=str(result)
                    )
                else:
                    yield result
    
    async def process_single_document(self, document: Dict[str, Any]) -> ProcessingResult:
        """
        Process a single document with comprehensive error handling and metrics
        """
        start_time = datetime.utcnow()
        doc_id = document.get('id', hashlib.md5(document['path'].encode()).hexdigest())
        
        try:
            self.logger.info(f"Processing document: {document['path']}")
            
            # Create document metadata
            metadata = await self._extract_document_metadata(document)
            
            # Download document content
            content = await self._download_document_async(document['download_url'])
            
            # Extract text content based on document type
            extracted_content = await self._extract_content_by_type(
                content, metadata.type, metadata.filename
            )
            
            # Create intelligent chunks
            chunks = await self._create_intelligent_chunks(extracted_content, metadata)
            
            # Store chunks with metadata
            if chunks:
                await self._store_chunks_async(chunks)
                
                # Calculate processing metrics
                end_time = datetime.utcnow()
                processing_time = int((end_time - start_time).total_seconds() * 1000)
                
                self.stats["documents_processed"] += 1
                self.stats["total_chunks_created"] += len(chunks)
                
                return ProcessingResult(
                    document_id=doc_id,
                    status=ProcessingStatus.COMPLETED,
                    chunks_created=len(chunks),
                    processing_time_ms=processing_time,
                    tokens_processed=sum(len(chunk.get('chunk', '').split()) for chunk in chunks),
                    confidence_score=await self._calculate_quality_score(chunks),
                    metadata=asdict(metadata)
                )
            else:
                return ProcessingResult(
                    document_id=doc_id,
                    status=ProcessingStatus.SKIPPED,
                    error_message="No extractable content found"
                )
                
        except Exception as e:
            self.logger.error(f"Error processing document {document['path']}: {str(e)}")
            self.stats["processing_errors"] += 1
            
            return ProcessingResult(
                document_id=doc_id,
                status=ProcessingStatus.FAILED,
                error_message=str(e)
            )
    
    async def _extract_document_metadata(self, document: Dict[str, Any]) -> DocumentMetadata:
        """Extract comprehensive document metadata"""
        
        # Determine document type
        extension = document.get('extension', '').lower()
        doc_type = self._get_document_type(extension)
        
        # Extract client information
        client_info = self.client_extractor.extract_client_info(document['path'])
        
        return DocumentMetadata(
            id=document.get('id', hashlib.md5(document['path'].encode()).hexdigest()),
            filename=document['name'],
            path=document['path'],
            size=document['size'],
            type=doc_type,
            processing_status=ProcessingStatus.PENDING,
            client_name=client_info['client_name'],
            pm_name=client_info['pm_name'],
            document_category=client_info['document_category'],
            is_client_specific=client_info['is_client_specific']
        )
    
    def _get_document_type(self, extension: str) -> DocumentType:
        """Determine document type from extension"""
        type_map = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.WORD,
            '.doc': DocumentType.WORD,
            '.xlsx': DocumentType.EXCEL,
            '.xls': DocumentType.EXCEL,
            '.pptx': DocumentType.POWERPOINT,
            '.ppt': DocumentType.POWERPOINT,
            '.txt': DocumentType.TEXT
        }
        return type_map.get(extension, DocumentType.UNKNOWN)
    
    async def _download_document_async(self, download_url: str) -> bytes:
        """Async document download with retry logic"""
        max_retries = 3
        retry_delay = 1
        
        for attempt in range(max_retries):
            try:
                async with aiohttp.ClientSession() as session:
                    async with session.get(download_url) as response:
                        if response.status == 200:
                            return await response.read()
                        else:
                            raise Exception(f"HTTP {response.status}: {response.reason}")
                            
            except Exception as e:
                if attempt == max_retries - 1:
                    raise e
                
                self.logger.warning(f"Download attempt {attempt + 1} failed, retrying: {e}")
                await asyncio.sleep(retry_delay * (2 ** attempt))
        
        raise Exception("Failed to download document after retries")
    
    async def _extract_content_by_type(self, 
                                     content: bytes, 
                                     doc_type: DocumentType, 
                                     filename: str) -> Any:
        """Extract content using appropriate method based on document type"""
        
        if doc_type == DocumentType.EXCEL:
            # Use enhanced Excel processor
            return await asyncio.get_event_loop().run_in_executor(
                self.executor,
                self.excel_processor.extract_from_excel,
                content,
                filename
            )
        
        elif doc_type in [DocumentType.PDF, DocumentType.WORD]:
            # Use Azure Document Intelligence
            return await self._extract_with_document_intelligence(content)
        
        elif doc_type == DocumentType.POWERPOINT:
            # Use PowerPoint processor
            return await asyncio.get_event_loop().run_in_executor(
                self.executor,
                self._extract_powerpoint_content,
                content
            )
        
        elif doc_type == DocumentType.TEXT:
            return content.decode('utf-8', errors='ignore')
        
        else:
            raise Exception(f"Unsupported document type: {doc_type}")
    
    async def _extract_with_document_intelligence(self, content: bytes) -> str:
        """Extract content using Azure Document Intelligence with async support"""
        try:
            credential = AzureKeyCredential(self.doc_intelligence_key)
            
            async with DocumentAnalysisClient(
                endpoint=self.doc_intelligence_endpoint,
                credential=credential
            ) as client:
                
                poller = await client.begin_analyze_document(
                    "prebuilt-read",
                    document=content
                )
                result = await poller.result()
                
                # Extract text content
                text_content = []
                for page in result.pages:
                    for line in page.lines:
                        text_content.append(line.content)
                
                return '\n'.join(text_content)
                
        except Exception as e:
            self.logger.error(f"Document Intelligence extraction failed: {e}")
            raise
    
    def _extract_powerpoint_content(self, content: bytes) -> str:
        """Extract content from PowerPoint files (runs in thread pool)"""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            prs = Presentation(BytesIO(content))
            all_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = [f"Slide {i+1}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                all_text.append('\n'.join(slide_text))
            
            return '\n\n'.join(all_text)
            
        except Exception as e:
            self.logger.error(f"PowerPoint extraction failed: {e}")
            return ""
    
    async def _create_intelligent_chunks(self, 
                                       content: Any, 
                                       metadata: DocumentMetadata) -> List[Dict[str, Any]]:
        """Create intelligent chunks with context preservation"""
        
        if metadata.type == DocumentType.EXCEL and isinstance(content, dict):
            # Handle Excel with table-aware chunking
            return await self._chunk_excel_content(content, metadata)
        
        elif isinstance(content, str):
            # Handle text content with semantic chunking
            return await self._chunk_text_content(content, metadata)
        
        else:
            self.logger.warning(f"Unknown content type for chunking: {type(content)}")
            return []
    
    async def _chunk_excel_content(self, 
                                 excel_data: Dict[str, Any], 
                                 metadata: DocumentMetadata) -> List[Dict[str, Any]]:
        """Create chunks for Excel content preserving table structures"""
        
        if excel_data.get("type") != "excel_sheets_enhanced":
            return []
        
        chunks = []
        
        for sheet_name, sheet_data in excel_data["sheets"].items():
            content = sheet_data["content"]
            tables = sheet_data["tables"]
            client_info = sheet_data.get("client_info")
            
            if not content or len(content.strip()) < 50:
                continue
            
            # Create sheet-level chunks
            chunk = {
                "chunk": content.strip(),
                "chunk_id": f"{metadata.id}_sheet_{sheet_name}",
                "document_id": metadata.id,
                "chunk_type": "excel_sheet_enhanced",
                "sheet_name": sheet_name,
                "table_count": len(tables),
                "table_types": [t.get("table_type") for t in tables],
                "chunk_length": len(content),
                "chunk_word_count": len(content.split()),
                "processed_timestamp": datetime.utcnow().isoformat(),
                
                # Document metadata
                "filename": metadata.filename,
                "document_path": metadata.path,
                "client_name": metadata.client_name,
                "pm_name": metadata.pm_name,
                "document_category": metadata.document_category,
                "is_client_specific": metadata.is_client_specific,
                
                # Processing metadata
                "processing_method": "enhanced_excel_processor",
                "confidence_score": 0.9  # High confidence for structured data
            }
            
            # Add sheet-level client info if detected
            if client_info:
                chunk.update({
                    "sheet_client_name": client_info.get("sheet_client_name"),
                    "sheet_client_confidence": client_info.get("confidence", 0.0)
                })
            
            chunks.append(chunk)
        
        return chunks
    
    async def _chunk_text_content(self, 
                                content: str, 
                                metadata: DocumentMetadata) -> List[Dict[str, Any]]:
        """Create semantic chunks for text content"""
        
        if not content or len(content.strip()) < 50:
            return []
        
        chunks = []
        content = content.strip()
        
        # Simple chunking with word boundaries (enhanced version would use semantic splitting)
        start = 0
        chunk_id = 0
        
        while start < len(content):
            end = min(start + self.chunk_size, len(content))
            
            # Find word boundary
            if end < len(content):
                while end > start and content[end] not in [' ', '\n', '\t', '.', '!', '?']:
                    end -= 1
                
                if end == start:
                    end = start + self.chunk_size
            
            chunk_text = content[start:end].strip()
            
            if chunk_text:
                chunk = {
                    "chunk": chunk_text,
                    "chunk_id": f"{metadata.id}_chunk_{chunk_id}",
                    "document_id": metadata.id,
                    "chunk_type": "text_semantic",
                    "chunk_index": chunk_id,
                    "chunk_length": len(chunk_text),
                    "chunk_word_count": len(chunk_text.split()),
                    "chunk_start_position": start,
                    "chunk_end_position": end,
                    "processed_timestamp": datetime.utcnow().isoformat(),
                    
                    # Document metadata
                    "filename": metadata.filename,
                    "document_path": metadata.path,
                    "client_name": metadata.client_name,
                    "pm_name": metadata.pm_name,
                    "document_category": metadata.document_category,
                    "is_client_specific": metadata.is_client_specific,
                    
                    # Processing metadata
                    "processing_method": "semantic_chunking",
                    "confidence_score": await self._calculate_chunk_quality(chunk_text)
                }
                
                chunks.append(chunk)
                chunk_id += 1
            
            # Move start position with overlap
            start = max(start + 1, end - self.chunk_overlap)
            
            if start >= len(content):
                break
        
        return chunks
    
    async def _calculate_chunk_quality(self, chunk_text: str) -> float:
        """Calculate quality score for a chunk"""
        
        # Simple quality metrics (enhanced version would use ML models)
        score = 0.5  # Base score
        
        # Length-based scoring
        if 100 <= len(chunk_text) <= 2000:
            score += 0.2
        
        # Sentence completion scoring
        if chunk_text.strip().endswith(('.', '!', '?')):
            score += 0.1
        
        # Word diversity scoring
        words = chunk_text.lower().split()
        unique_words = len(set(words))
        if len(words) > 0:
            diversity = unique_words / len(words)
            score += diversity * 0.2
        
        return min(1.0, score)
    
    async def _calculate_quality_score(self, chunks: List[Dict[str, Any]]) -> float:
        """Calculate overall quality score for processed chunks"""
        if not chunks:
            return 0.0
        
        scores = [chunk.get('confidence_score', 0.5) for chunk in chunks]
        return sum(scores) / len(scores)
    
    async def _store_chunks_async(self, chunks: List[Dict[str, Any]]) -> None:
        """Store chunks in Azure Blob Storage with async operations"""
        try:
            container_name = "jennifur-processed"
            
            async with BlobServiceClient.from_connection_string(
                self.storage_connection
            ) as blob_service:
                
                container_client = blob_service.get_container_client(container_name)
                
                # Ensure container exists
                try:
                    await container_client.get_container_properties()
                except:
                    await container_client.create_container()
                
                # Store chunks concurrently
                async def store_single_chunk(chunk):
                    chunk_id = chunk.get("chunk_id", "unknown")
                    blob_name = f"{chunk.get('document_id', 'unknown')}/{chunk_id}.json"
                    
                    blob_client = container_client.get_blob_client(blob_name)
                    chunk_json = json.dumps(chunk, indent=2)
                    
                    await blob_client.upload_blob(
                        chunk_json,
                        overwrite=True,
                        content_type="application/json"
                    )
                
                # Store all chunks concurrently
                tasks = [store_single_chunk(chunk) for chunk in chunks]
                await asyncio.gather(*tasks)
                
                self.logger.info(f"Stored {len(chunks)} chunks in blob storage")
                
        except Exception as e:
            self.logger.error(f"Error storing chunks: {e}")
            raise
    
    async def get_processing_statistics(self) -> Dict[str, Any]:
        """Get current processing statistics"""
        runtime = (datetime.utcnow() - self.stats["start_time"]).total_seconds()
        
        return {
            "documents_processed": self.stats["documents_processed"],
            "total_chunks_created": self.stats["total_chunks_created"],
            "processing_errors": self.stats["processing_errors"],
            "runtime_seconds": runtime,
            "documents_per_minute": (self.stats["documents_processed"] / runtime * 60) if runtime > 0 else 0,
            "chunks_per_document": (self.stats["total_chunks_created"] / self.stats["documents_processed"]) if self.stats["documents_processed"] > 0 else 0
        }
    
    async def close(self):
        """Clean up resources"""
        self.executor.shutdown(wait=True)
        self.logger.info("Enhanced Document Processor closed")


# Agentic Processing Extensions
class AgenticProcessingOrchestrator:
    """
    Orchestrates intelligent document processing with agent-like capabilities
    """
    
    def __init__(self, processor: EnhancedDocumentProcessor):
        self.processor = processor
        self.logger = logging.getLogger(__name__)
    
    async def intelligent_batch_processing(self, 
                                         documents: List[Dict[str, Any]],
                                         priority_strategy: str = "size_first") -> AsyncIterator[ProcessingResult]:
        """
        Process documents with intelligent prioritization and resource management
        """
        
        # Sort documents by priority strategy
        if priority_strategy == "size_first":
            documents.sort(key=lambda x: x.get('size', 0))
        elif priority_strategy == "client_priority":
            documents.sort(key=lambda x: self._get_client_priority(x.get('path', '')))
        
        # Process with adaptive batch sizing
        batch_size = await self._calculate_optimal_batch_size(documents)
        
        self.logger.info(f"Processing {len(documents)} documents with batch size {batch_size}")
        
        async for result in self.processor.process_documents_batch(documents, batch_size):
            # Adaptive processing based on results
            if result.status == ProcessingStatus.FAILED:
                self.logger.warning(f"Document processing failed: {result.error_message}")
            
            yield result
    
    def _get_client_priority(self, document_path: str) -> int:
        """Assign priority based on client importance (0 = highest priority)"""
        high_priority_clients = ['important_client_1', 'vip_client']
        
        for i, client in enumerate(high_priority_clients):
            if client.lower() in document_path.lower():
                return i
        
        return 999  # Default low priority
    
    async def _calculate_optimal_batch_size(self, documents: List[Dict[str, Any]]) -> int:
        """Calculate optimal batch size based on document characteristics"""
        
        if not documents:
            return 1
        
        # Simple heuristic based on average document size
        avg_size = sum(doc.get('size', 0) for doc in documents) / len(documents)
        
        if avg_size < 1024 * 1024:  # < 1MB
            return 20
        elif avg_size < 10 * 1024 * 1024:  # < 10MB
            return 10
        else:
            return 5